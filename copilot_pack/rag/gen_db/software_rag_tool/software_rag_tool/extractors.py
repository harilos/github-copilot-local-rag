from __future__ import annotations

import platform
import shutil
import subprocess
import tempfile
from pathlib import Path

from .chunking import TextSection, chunk_text, normalize_text

SUPPORTED_EXTENSIONS = {
    ".md",
    ".txt",
    ".log",
    ".pdf",
    ".docx",
    ".doc",
    ".pptx",
    ".ppt",
    ".xlsx",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".java",
    ".go",
    ".rs",
    ".cs",
    ".rb",
    ".php",
    ".sh",
    ".ps1",
    ".sql",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
}


def extract_sections(path: Path) -> list[TextSection]:
    ext = path.suffix.lower()
    if ext in {".md", ".txt", ".log", ".py", ".js", ".jsx", ".ts", ".tsx", ".java", ".go", ".rs", ".cs", ".rb", ".php", ".sh", ".ps1", ".sql", ".json", ".yaml", ".yml", ".toml", ".ini"}:
        return _extract_plain(path)
    if ext == ".pdf":
        return _extract_pdf(path)
    if ext == ".docx":
        return _extract_docx(path)
    if ext == ".pptx":
        return _extract_pptx(path)
    if ext == ".xlsx":
        return _extract_xlsx(path)
    if ext in {".doc", ".ppt"}:
        return _extract_legacy_office(path)
    return []


def _extract_plain(path: Path) -> list[TextSection]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return chunk_text(path.name, text)


def _extract_pdf(path: Path) -> list[TextSection]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    sections: list[TextSection] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        sections.extend(chunk_text(f"Page {i}", text))
    if sections:
        return sections
    return [TextSection(title=path.name, text="")]


def _extract_docx(path: Path) -> list[TextSection]:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return chunk_text(path.name, "\n".join(parts))


def _extract_pptx(path: Path) -> list[TextSection]:
    from pptx import Presentation

    prs = Presentation(str(path))
    sections: list[TextSection] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        sections.extend(chunk_text(f"Slide {i}", "\n".join(parts)))
    return sections


def _extract_xlsx(path: Path) -> list[TextSection]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sections: list[TextSection] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(v) for v in row if v is not None and str(v).strip()]
            if values:
                rows.append("\t".join(values))
        sections.extend(chunk_text(sheet.title, "\n".join(rows)))
    return sections


def _extract_legacy_office(path: Path) -> list[TextSection]:
    if path.suffix.lower() == ".doc" and platform.system() == "Darwin":
        text = _convert_doc_with_textutil(path)
        if text:
            return chunk_text(path.name, text)

    text = _convert_with_libreoffice(path)
    if text:
        return chunk_text(path.name, text)
    raise RuntimeError(f"Failed to extract legacy Office file: {path}")


def _convert_doc_with_textutil(path: Path) -> str:
    exe = shutil.which("textutil")
    if not exe:
        return ""
    proc = subprocess.run(
        [exe, "-convert", "txt", "-stdout", str(path)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        timeout=60,
    )
    if proc.returncode != 0:
        return ""
    return normalize_text(proc.stdout)


def _convert_with_libreoffice(path: Path) -> str:
    exe = shutil.which("soffice") or shutil.which("libreoffice")
    if not exe:
        return ""
    with tempfile.TemporaryDirectory(prefix="ac-rag-office-") as tmp:
        tmp_path = Path(tmp)
        profile = tmp_path / "profile"
        cmd = [
            exe,
            f"-env:UserInstallation=file://{profile}",
            "--headless",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(tmp_path),
            str(path),
        ]
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=90)
        if proc.returncode != 0:
            return ""
        txt_files = sorted(tmp_path.glob("*.txt"))
        if not txt_files:
            return ""
        return normalize_text(txt_files[0].read_text(encoding="utf-8", errors="replace"))
