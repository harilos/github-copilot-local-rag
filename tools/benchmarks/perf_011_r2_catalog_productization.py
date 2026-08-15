from __future__ import annotations

import argparse
import hashlib
import json
import math
import os
import random
import shutil
import statistics
import subprocess
import sys
import tempfile
import time
import zipfile
from contextlib import contextmanager
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Iterator


SEED = 1102
DOCX_COUNT = 40
PPTX_COUNT = 20
DOCX_RECORDS = 5
PPTX_RECORDS = 10
TOTAL_RECORDS = 400
UPDATED_DOCX = tuple(range(4))
UPDATED_PPTX = tuple(range(2))
UPDATED_RECORDS = 40


def _sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _canonical_hash(value: Any) -> str:
    raw = json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":"))
    return _sha256_bytes(raw.encode("utf-8"))


def _canonicalize_openxml(path: Path) -> None:
    temporary = path.with_suffix(path.suffix + ".canonical")
    with zipfile.ZipFile(path) as source, zipfile.ZipFile(
        temporary,
        "w",
        compression=zipfile.ZIP_DEFLATED,
        compresslevel=9,
    ) as destination:
        for name in sorted(source.namelist()):
            original = source.getinfo(name)
            info = zipfile.ZipInfo(name, date_time=(1980, 1, 1, 0, 0, 0))
            info.compress_type = zipfile.ZIP_DEFLATED
            info.create_system = original.create_system
            info.external_attr = original.external_attr
            destination.writestr(info, source.read(name))
    os.replace(temporary, path)


def _generate_corpus(initial: Path, updated: Path) -> dict[str, Any]:
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches

    initial.mkdir(parents=True)
    updated.mkdir(parents=True)
    generator = random.Random(SEED)

    def docx(path: Path, index: int, generation: int, fixture_token: int) -> None:
        document = Document()
        document.core_properties.author = "Local RAG PERF-011-R2"
        document.core_properties.created = datetime(2026, 8, 15, tzinfo=timezone.utc)
        document.core_properties.modified = datetime(2026, 8, 15, tzinfo=timezone.utc)
        document.add_heading(f"Catalog operations guide {index:02d}", level=0)
        words = (
            "catalog office update alpha beta gamma æ—¥æœ¬èªž project "
            f"DOCX{index:02d} generation{generation} seed{fixture_token} RFC{110000 + index}"
        )
        for section in range(DOCX_RECORDS):
            document.add_heading(f"Section {section + 1}", level=1)
            body = " ".join(f"{words} item{section:02d}-{repeat:02d}." for repeat in range(4))
            document.add_paragraph(body)
        table = document.add_table(rows=2, cols=3)
        for row in range(2):
            for column in range(3):
                table.cell(row, column).text = (
                    f"Metric {index:02d}-{row}-{column} generation {generation}"
                )
        document.save(path)
        _canonicalize_openxml(path)

    def pptx(path: Path, index: int, generation: int, fixture_token: int) -> None:
        presentation = Presentation()
        presentation.core_properties.author = "Local RAG PERF-011-R2"
        presentation.core_properties.created = datetime(2026, 8, 15, tzinfo=timezone.utc)
        presentation.core_properties.modified = datetime(2026, 8, 15, tzinfo=timezone.utc)
        presentation.slide_width = Inches(13.333)
        presentation.slide_height = Inches(7.5)
        for slide_index in range(PPTX_RECORDS):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.7))
            title.text_frame.text = f"Quarterly catalog plan {index:02d} / {slide_index + 1}"
            body = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.5), Inches(4.8))
            frame = body.text_frame
            frame.text = (
                f"PPTX{index:02d} slide {slide_index + 1} generation {generation} "
                f"seed {fixture_token} catalog update æ—¥æœ¬èªž RFC{120000 + index}."
            )
            for bullet in range(3):
                paragraph = frame.add_paragraph()
                paragraph.level = 0
                paragraph.text = (
                    f"Action {bullet + 1}: validate Office record {index:02d}-{slide_index:02d}."
                )
        presentation.save(path)
        _canonicalize_openxml(path)

    for index in range(DOCX_COUNT):
        name = f"word/catalog-guide-{index:02d}.docx"
        target = initial / name
        target.parent.mkdir(parents=True, exist_ok=True)
        docx(target, index, 0, generator.randrange(1_000_000_000))
    for index in range(PPTX_COUNT):
        name = f"powerpoint/catalog-plan-{index:02d}.pptx"
        target = initial / name
        target.parent.mkdir(parents=True, exist_ok=True)
        pptx(target, index, 0, generator.randrange(1_000_000_000))
    shutil.copytree(initial, updated, dirs_exist_ok=True)
    for index in UPDATED_DOCX:
        docx(
            updated / f"word/catalog-guide-{index:02d}.docx",
            index,
            1,
            generator.randrange(1_000_000_000),
        )
    for index in UPDATED_PPTX:
        pptx(
            updated / f"powerpoint/catalog-plan-{index:02d}.pptx",
            index,
            1,
            generator.randrange(1_000_000_000),
        )

    def hashes(root: Path) -> list[dict[str, Any]]:
        return [
            {
                "path": path.relative_to(root).as_posix(),
                "bytes": path.stat().st_size,
                "sha256": _sha256_bytes(path.read_bytes()),
            }
            for path in sorted(root.rglob("*"))
            if path.is_file()
        ]

    initial_hashes = hashes(initial)
    updated_hashes = hashes(updated)
    return {
        "seed": SEED,
        "initial_files": initial_hashes,
        "updated_files": updated_hashes,
        "initial_file_manifest_hash": _canonical_hash(initial_hashes),
        "updated_file_manifest_hash": _canonical_hash(updated_hashes),
    }


def _configure(code_root: Path, output_root: Path, model_dir: Path) -> None:
    package = code_root / ".copilot" / "rag" / "gen_db" / "software_rag_tool"
    sys.path.insert(0, str(package))
    os.environ.update(
        {
            "RAG_OUTPUT_ROOT": str(output_root),
            "CHROMA_DIR_V2": str(output_root / "index" / "chroma"),
            "CHROMA_COLLECTION": "perf011_r2",
            "EMBEDDING_BACKEND": "onnx",
            "EMBEDDING_MODEL": "cl-nagoya/ruri-v3-30m",
            "EMBEDDING_DIMENSION": "256",
            "EMBEDDING_QUANTIZATION": "dynamic-int8",
            "EMBEDDING_ONNX_DIR": str(model_dir),
            "LOCAL_RAG_LEXICAL_TOKENIZER": "sudachi",
            "EMBED_BATCH_SIZE": "8",
            "RAG_ONNX_THREADS": "4",
            "OMP_NUM_THREADS": "4",
            "TOKENIZERS_PARALLELISM": "false",
            "HF_HUB_OFFLINE": "1",
            "TRANSFORMERS_OFFLINE": "1",
        }
    )


def _ordered_records(corpus: Path) -> tuple[list[dict[str, Any]], dict[str, Any]]:
    from software_rag_tool.embeddings import get_document_token_budget
    from software_rag_tool.records import build_records_for_file, iter_input_files

    budget = get_document_token_budget()
    records: list[dict[str, Any]] = []
    by_file: dict[str, int] = {}
    token_manifest: list[dict[str, Any]] = []
    for path in iter_input_files(corpus):
        relative = path.relative_to(corpus).as_posix()
        built = build_records_for_file(corpus, path, source_id="perf011-r2")
        by_file[relative] = len(built)
        for item in built:
            token_ids = budget.tokenizer(
                budget.document_prefix + str(item["embedding_text"]),
                add_special_tokens=True,
            )["input_ids"]
            token_manifest.append(
                {
                    "path": relative,
                    "source_id": item["metadata"]["source_id"],
                    "id": item["id"],
                    "text": item["text"],
                    "metadata": item["metadata"],
                    "token_count": len(token_ids),
                    "token_ids": token_ids,
                }
            )
        records.extend(built)
    docx_counts = [count for path, count in by_file.items() if path.endswith(".docx")]
    pptx_counts = [count for path, count in by_file.items() if path.endswith(".pptx")]
    if len(docx_counts) != DOCX_COUNT or set(docx_counts) != {DOCX_RECORDS}:
        raise RuntimeError(f"DOCX fixture count mismatch: {docx_counts}")
    if len(pptx_counts) != PPTX_COUNT or set(pptx_counts) != {PPTX_RECORDS}:
        raise RuntimeError(f"PPTX fixture count mismatch: {pptx_counts}")
    if len(records) != TOTAL_RECORDS:
        raise RuntimeError(f"fixture record count mismatch: {len(records)}")
    return records, {
        "record_count": len(records),
        "docx_records": sum(docx_counts),
        "pptx_records": sum(pptx_counts),
        "per_file_counts": by_file,
        "ordered_record_manifest_hash": _canonical_hash(token_manifest),
    }


def _tree_bytes(root: Path) -> int:
    return sum(path.stat().st_size for path in root.rglob("*") if path.is_file())


def _peak_working_set() -> int:
    if os.name != "nt":
        return 0
    import ctypes
    from ctypes import wintypes

    class Counters(ctypes.Structure):
        _fields_ = [
            ("cb", wintypes.DWORD),
            ("PageFaultCount", wintypes.DWORD),
            ("PeakWorkingSetSize", ctypes.c_size_t),
            ("WorkingSetSize", ctypes.c_size_t),
            ("QuotaPeakPagedPoolUsage", ctypes.c_size_t),
            ("QuotaPagedPoolUsage", ctypes.c_size_t),
            ("QuotaPeakNonPagedPoolUsage", ctypes.c_size_t),
            ("QuotaNonPagedPoolUsage", ctypes.c_size_t),
            ("PagefileUsage", ctypes.c_size_t),
            ("PeakPagefileUsage", ctypes.c_size_t),
        ]

    counters = Counters()
    counters.cb = ctypes.sizeof(counters)
    psapi = ctypes.WinDLL("psapi", use_last_error=True)
    kernel32 = ctypes.WinDLL("kernel32", use_last_error=True)
    kernel32.GetCurrentProcess.restype = wintypes.HANDLE
    psapi.GetProcessMemoryInfo.argtypes = [
        wintypes.HANDLE,
        ctypes.POINTER(Counters),
        wintypes.DWORD,
    ]
    psapi.GetProcessMemoryInfo.restype = wintypes.BOOL
    handle = kernel32.GetCurrentProcess()
    if not psapi.GetProcessMemoryInfo(handle, ctypes.byref(counters), counters.cb):
        return 0
    return int(counters.PeakWorkingSetSize)


@contextmanager
def _catalog_trace(
    catalog: Any,
    samples: dict[str, Any],
    *,
    enabled: bool,
) -> Iterator[None]:
    if not enabled:
        yield
        return
    original = catalog.connect

    def traced(path: Path | None = None):
        manager = original(path)

        class Context:
            def __enter__(self):
                connection = manager.__enter__()
                self.connection = connection
                self.before = connection.total_changes
                if enabled:
                    connection.set_trace_callback(samples["sql"].append)
                return connection

            def __exit__(self, *args: Any):
                samples["total_changes"] += self.connection.total_changes - self.before
                if enabled:
                    self.connection.set_trace_callback(None)
                return manager.__exit__(*args)

        return Context()

    catalog.connect = traced
    try:
        yield
    finally:
        catalog.connect = original


def _timed_call(target: Any, name: str, totals: dict[str, float]) -> tuple[Any, Callable[..., Any]]:
    original = getattr(target, name)

    def wrapped(*args: Any, **kwargs: Any) -> Any:
        started = time.perf_counter()
        try:
            return original(*args, **kwargs)
        finally:
            totals[name] = totals.get(name, 0.0) + time.perf_counter() - started

    setattr(target, name, wrapped)
    return target, original


def _worker(args: argparse.Namespace) -> int:
    code_root = Path(args.code_root).resolve()
    output_root = Path(args.output_root).resolve()
    corpus = Path(args.corpus).resolve()
    model_dir = Path(args.model_dir).resolve()
    _configure(code_root, output_root, model_dir)
    records, manifest = _ordered_records(corpus)
    expected = json.loads(Path(args.expected_manifest).read_text(encoding="utf-8"))
    if manifest["ordered_record_manifest_hash"] != expected["initial"]["ordered_record_manifest_hash"]:
        raise RuntimeError("baseline/candidate record manifest mismatch")
    if args.cohort == "prepare":
        updated_records, updated_manifest = _ordered_records(Path(args.updated_corpus))
        changed_paths = {
            *(f"word/catalog-guide-{index:02d}.docx" for index in UPDATED_DOCX),
            *(f"powerpoint/catalog-plan-{index:02d}.pptx" for index in UPDATED_PPTX),
        }
        changed = [item for item in updated_records if item["metadata"]["path"].split("/", 1)[-1] in changed_paths]
        old_ids = [item["id"] for item in records if item["metadata"]["path"].split("/", 1)[-1] in changed_paths]
        if len(changed) != UPDATED_RECORDS or len(old_ids) != UPDATED_RECORDS:
            raise RuntimeError(f"update record mismatch: new={len(changed)} old={len(old_ids)}")
        Path(args.records).write_text(json.dumps({"initial": records, "updated": changed, "old_ids": old_ids}, ensure_ascii=False), encoding="utf-8")
        Path(args.output).write_text(json.dumps({"initial": manifest, "updated": updated_manifest}, ensure_ascii=False, indent=2), encoding="utf-8")
        return 0

    from software_rag_tool import catalog, incremental, manifest as product_manifest, store

    prepared = json.loads(Path(args.records).read_text(encoding="utf-8"))
    sql: dict[str, Any] = {"sql": [], "total_changes": 0}
    phases: dict[str, float] = {}
    restored: list[tuple[Any, str, Callable[..., Any]]] = []
    catalog_target = store if args.cohort == "core" else incremental
    vector_target = store if args.cohort == "core" else incremental
    for target, name in ((catalog_target, "upsert_catalog_records"), (vector_target, "upsert_records")):
        owner, original = _timed_call(target, name, phases)
        restored.append((owner, name, original))
    if args.cohort == "e2e":
        owner, original = _timed_call(incremental, "build_records_for_file", phases)
        restored.append((owner, "build_records_for_file", original))

    wall_start = time.perf_counter()
    cpu_start = time.process_time()
    try:
        with _catalog_trace(catalog, sql, enabled=args.trace_probe):
            if args.cohort == "core":
                from software_rag_tool.jsonl import write_jsonl
                from software_rag_tool.paths import clean_dir

                write_jsonl(clean_dir() / "perf011-r2.jsonl", prepared["initial"])
              ×ou¶‰žËkºwµçy‘ÌˆèÁ¡…Í•Ì¹•Ð ‰‰Õ¥±‘}É•½É‘Í}™½É}™¥±”ˆ°€À¸À¤°(€€€€€€€€‰ÍÅ±}½Õ¹Ðˆè±•¸¡ÍÅ±l‰ÍÅ°‰t¤°(€€€€€€€€‰Ñ½Ñ…±}¡…¹•ÌˆèÍÅ±l‰Ñ½Ñ…±}¡…¹•Ì‰t°(€€€€€€€€‰Á•…­}ÉÍÍ}‰åÑ•Ìˆè}Á•…­}Ý½É­¥¹}Í•Ð ¤°(€€€€€€€€‰‘‰}‰åÑ•Ìˆè}ÑÉ••}‰åÑ•Ì¡½ÕÑÁÕÑ}É½½Ð€¼€‰¥¹‘•àˆ¤°(€€€€€€€€‰Ý…±}‰åÑ•ÌˆèÍÕ´¡Á…Ñ ¹ÍÑ…Ð ¤¹ÍÑ}Í¥é”™½ÈÁ…Ñ ¥¸½ÕÑÁÕÑ}É½½Ð¹É±½ˆ ˆ¨µÝ…°ˆ¤¤°(€€€€€€€€‰Ù•Ñ½É}½Õ¹ÐˆèÍÑ½É”¹½±±•Ñ¥½¹}½Õ¹Ð ¤°(€€€€€€€€‰…Ñ…±½}½Õ¹Ðˆè…Ñ…±½}½Õ¹Ð°(€€€€€€€€‰µ…¹¥™•ÍÑ}É•½É‘}½Õ¹Ðˆè¥¹Ð¡µ…¹¥™•ÍÑ}Á…å±½…¹•Ð ‰É•½É‘}½Õ¹Ðˆ°€´Ä¤¤°(€€€€€€€€‰µ…¹¥™•ÍÐˆèµ…¹¥™•ÍÑ}Á…å±½…°(€€€€€€€€‰•µ‰•‘‘•É}±…ÍÌˆèÑåÁ”¡•µ‰•‘‘•È¤¹}}¹…µ•}|°(€€€€€€€€‰½¹¹á}ÁÉ½Ù¥‘•ÉÌˆèÁÉ½Ù¥‘•ÉÌ°(€€€€€€€€‰½±±•Ñ¥½¹}¹…µ”ˆèÍÑ½É”¹½±±•Ñ¥½¹}¹…µ” ¤°(€€€€€€€€‰½±±•Ñ¥½¹}µ•Ñ…‘…Ñ„ˆè‘¥Ð¡½±±•Ñ¥½¸¹µ•Ñ…‘…Ñ„½Èíô¤¥˜½±±•Ñ¥½¸¥Ì¹½Ð9½¹”•±Í”íô°(€€€€€€€€‰¡É½µ…}ÍÅ±¥Ñ•}‰åÑ•Ìˆè¡É½µ…}ÍÅ±¥Ñ”¹ÍÑ…Ð ¤¹ÍÑ}Í¥é”¥˜¡É½µ…}ÍÅ±¥Ñ”¹•á¥ÍÑÌ ¤•±Í”€À°(€€€€€€€€‰¥¹Ñ•É¥Ñäˆè¥¹Ñ•É¥Ñä°(€€€€€€€€‰™½É•¥¹}­•å}•ÉÉ½ÉÌˆè™½É•¥¹}­•åÌ°(€€€€€€€€‰É•½É‘}µ…¹¥™•ÍÑ}¡…Í ˆèµ…¹¥™•ÍÑl‰½É‘•É•‘}É•½É‘}µ…¹¥™•ÍÑ}¡…Í ‰t°(€€€€€€€€‰ÑÉ…•}ÁÉ½‰”ˆè‰½½°¡…ÉÌ¹ÑÉ…•}ÁÉ½‰”¤°(€€€ô(€€€¥˜¹½Ð€ (€€€€€€€É•ÍÕ±Ñl‰Ù•Ñ½É}½Õ¹Ð‰t€ôôQ=Q1}I=IL(€€€€€€€…¹É•ÍÕ±Ñl‰…Ñ…±½}½Õ¹Ð‰t€ôôQ=Q1}I=IL(€€€€€€€…¹É•ÍÕ±Ñl‰µ…¹¥™•ÍÑ}É•½É‘}½Õ¹Ð‰t€ôôQ=Q1}I=IL(€€€€€€€…¹¥¹Ñ•É¥Ñä€ôô€‰½¬ˆ(€€€€€€€…¹™½É•¥¹}­•åÌ€ôô€À(€€€€€€€…¹ÑåÁ”¡•µ‰•‘‘•È¤¹}}¹…µ•}|€ôô€‰=¹¹áIÕ¹Ñ¥µ•µ‰•‘‘•Èˆ(€€€€€€€…¹€‰AUá•ÕÑ¥½¹AÉ½Ù¥‘•Èˆ¥¸ÁÉ½Ù¥‘•ÉÌ(€€€€€€€…¹¡É½µ…}ÍÅ±¥Ñ”¹¥Í}™¥±” ¤(€€€€€€€…¹µ…¹¥™•ÍÑ}Á…å±½…¹•Ð ‰½±±•Ñ¥½¸ˆ¤€ôô€‰Á•É˜ÀÄÅ}ÈÈˆ(€€€€€€€…¹µ…¹¥™•ÍÑ}Á…å±½…¹•Ð ‰•µ‰•‘‘¥¹}µ½‘•°ˆ¤€ôô€‰°µ¹…½å„½ÉÕÉ¤µØÌ´ÌÁ´ˆ(€€€€€€€…¹µ…¹¥™•ÍÑ}Á…å±½…¹•Ð ‰•µ‰•‘‘¥¹}‰…­•¹ˆ¤€ôô€‰½¹¹àˆ(€€€€€€€…¹µ…¹¥™•ÍÑ}Á…å±½…¹•Ð ‰•µ‰•‘‘¥¹}‘¥µ•¹Í¥½¸ˆ¤€ôô€ÈÔØ(€€€€€€€…¹µ…¹¥™•ÍÑ}Á…å±½…¹•Ð ‰ÅÕ…¹Ñ¥é…Ñ¥½¸ˆ¤€ôô€‰‘å¹…µ¥Œµ¥¹Ðàˆ(€€€€€€€…¹½±±•Ñ¥½¸¥Ì¹½Ð9½¹”(€€€€€€€…¹½±±•Ñ¥½¸¹µ•Ñ…‘…Ñ„¹•Ð ‰•µ‰•‘‘¥¹}µ½‘•°ˆ¤€ôô€‰°µ¹…½å„½ÉÕÉ¤µØÌ´ÌÁ´ˆ(€€€€€€€…¹½±±•Ñ¥½¸¹µ•Ñ…‘…Ñ„¹•Ð ‰•µ‰•‘‘¥¹}‰…­•¹ˆ¤€ôô€‰½¹¹àˆ(€€€€€€€…¹½±±•Ñ¥½¸¹µ•Ñ…‘…Ñ„¹•Ð ‰•µ‰•‘‘¥¹}‘¥µ•¹Í¥½¸ˆ¤€ôô€ÈÔØ(€€€€€€€…¹½±±•Ñ¥½¸¹µ•Ñ…‘…Ñ„¹•Ð ‰ÅÕ…¹Ñ¥é…Ñ¥½¸ˆ¤€ôô€‰‘å¹…µ¥Œµ¥¹Ðàˆ(€€€€¤è(€€€€€€€É…¥Í”IÕ¹Ñ¥µ•ÉÉ½È¡˜‰ÁÉ½‘ÕÐ½±µ¥Íµ…Ñ èíÉ•ÍÕ±Ñôˆ¤(€€€A…Ñ ¡…ÉÌ¹½ÕÑÁÕÐ¤¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡É•ÍÕ±Ð°•¹ÍÕÉ•}…Í¥¤õ…±Í”°¥¹‘•¹ÐôÈ¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€É•ÑÕÉ¸€À(()‘•˜}ÀäÔ¡Ù…±Õ•Ìè±¥ÍÑm™±½…Ñt¤€´ø™±½…Ðè(€€€É•ÑÕÉ¸Í½ÉÑ•¡Ù…±Õ•Ì¥mµ…à À°µ…Ñ ¹•¥°¡±•¸¡Ù…±Õ•Ì¤€¨€À¸äÔ¤€´€Ä¥t(()‘•˜}ÍÕµµ…Éä¡Í…µÁ±•Ìè±¥ÍÑm‘¥ÑmÍÑÈ°¹åut¤€´ø‘¥ÑmÍÑÈ°¹åtè(€€€É•ÍÕ±Ðè‘¥ÑmÍÑÈ°¹åt€ôì‰ÉÕ¹Ìˆè±•¸¡Í…µÁ±•Ì¥ô(€€€¹Õµ•É¥Œ€ôm­•ä™½È­•ä°Ù…±Õ”¥¸Í…µÁ±•ÍlÁt¹¥Ñ•µÌ ¤¥˜¥Í¥¹ÍÑ…¹”¡Ù…±Õ”°€¡¥¹Ð°™±½…Ð¤¤…¹¹½Ð¥Í¥¹ÍÑ…¹”¡Ù…±Õ”°‰½½°¥t(€€€™½È­•ä¥¸¹Õµ•É¥Œè(€€€€€€€Ù…±Õ•Ì€ôm™±½…Ð¡Í…µÁ±•m­•åt¤™½ÈÍ…µÁ±”¥¸Í…µÁ±•Ít(€€€€€€€É•ÍÕ±Ñm˜‰í­•åõ}ÀÔÀ‰t€ôÍÑ…Ñ¥ÍÑ¥Ì¹µ•‘¥…¸¡Ù…±Õ•Ì¤(€€€€€€€É•ÍÕ±Ñm˜‰í­•åõ}ÀäÔ‰t€ô}ÀäÔ¡Ù…±Õ•Ì¤(€€€É•ÑÕÉ¸É•ÍÕ±Ð(()‘•˜}É•‘ÕÑ¥½¸¡‰…Í•±¥¹”è±¥ÍÑm‘¥ÑmÍÑÈ°¹åut°…¹‘¥‘…Ñ”è±¥ÍÑm‘¥ÑmÍÑÈ°¹åut°­•äèÍÑÈ¤€´ø™±½…Ðè(€€€É…Ñ¥½Ì€ôlÄ¸À€´™±½…Ð¡m­•åt¤€¼™±½…Ð¡‰m­•åt¤™½Èˆ°Œ¥¸é¥À¡‰…Í•±¥¹”°…¹‘¥‘…Ñ”°ÍÑÉ¥ÐõQÉÕ”¥t(€€€É•ÑÕÉ¸ÍÑ…Ñ¥ÍÑ¥Ì¹µ•‘¥…¸¡É…Ñ¥½Ì¤(()‘•˜}ÉÕ¸¡…ÉÌè…ÉÁ…ÉÍ”¹9…µ•ÍÁ…”¤€´ø¥¹Ðè(€€€ÍÉ¥ÁÐ€ôA…Ñ ¡}}™¥±•}|¤¹É•Í½±Ù” ¤(€€€½ÕÑÁÕÐ€ôA…Ñ ¡…ÉÌ¹½ÕÑÁÕÐ¤¹É•Í½±Ù” ¤(€€€½ÕÑÁÕÐ¹Á…É•¹Ð¹µ­‘¥È¡Á…É•¹ÑÌõQÉÕ”°•á¥ÍÑ}½¬õQÉÕ”¤(€€€ÍÑ…ÉÑ•€ôÑ¥µ”¹µ½¹½Ñ½¹¥Œ ¤(€€€Ý¥Ñ Ñ•µÁ™¥±”¹Q•µÁ½É…Éå¥É•Ñ½Éä¡ÁÉ•™¥àô‰Á•É˜ÀÄÄµÈÈ´ˆ¤…ÌÑ•µÁ½É…Éäè(€€€€€€€É½½Ð€ôA…Ñ ¡Ñ•µÁ½É…Éä¤(€€€€€€€¥¹¥Ñ¥…°€ôÉ½½Ð€¼€‰™¥áÑÕÉ”ˆ€¼€‰¥¹¥Ñ¥…°ˆ€¼€‰Í½ÕÉ”ˆ(€€€€€€€ÕÁ‘…Ñ•€ôÉ½½Ð€¼€‰™¥áÑÕÉ”ˆ€¼€‰ÕÁ‘…Ñ•ˆ€¼€‰Í½ÕÉ”ˆ(€€€€€€€™¥áÑÕÉ”€ô}•¹•É…Ñ•}½ÉÁÕÌ¡¥¹¥Ñ¥…°°ÕÁ‘…Ñ•¤(€€€€€€€É•ÁÉ½‘Õ¥‰±•}¥¹¥Ñ¥…°€ôÉ½½Ð€¼€‰™¥áÑÕÉ”µÉ•ÁÉ¼ˆ€¼€‰¥¹¥Ñ¥…°ˆ€¼€‰Í½ÕÉ”ˆ(€€€€€€€É•ÁÉ½‘Õ¥‰±•}ÕÁ‘…Ñ•€ôÉ½½Ð€¼€‰™¥áÑÕÉ”µÉ•ÁÉ¼ˆ€¼€‰ÕÁ‘…Ñ•ˆ€¼€‰Í½ÕÉ”ˆ(€€€€€€€É•ÁÉ½‘Õ•€ô}•¹•É…Ñ•}½ÉÁÕÌ¡É•ÁÉ½‘Õ¥‰±•}¥¹¥Ñ¥…°°É•ÁÉ½‘Õ¥‰±•}ÕÁ‘…Ñ•¤(€€€€€€€¥˜€ (€€€€€€€€€€€™¥áÑÕÉ•l‰¥¹¥Ñ¥…±}™¥±•}µ…¹¥™•ÍÑ}¡…Í ‰t€„ôÉ•ÁÉ½‘Õ•‘l‰¥¹¥Ñ¥…±}™¥±•}µ…¹¥™•ÍÑ}¡…Í ‰t(€€€€€€€€€€€½È™¥áÑÕÉ•l‰ÕÁ‘…Ñ•‘}™¥±•}µ…¹¥™•ÍÑ}¡…Í ‰t€„ôÉ•ÁÉ½‘Õ•‘l‰ÕÁ‘…Ñ•‘}™¥±•}µ…¹¥™•ÍÑ}¡…Í ‰t(€€€€€€€€¤è(€€€€€€€€€€€É…¥Í”IÕ¹Ñ¥µ•ÉÉ½È ‰=™™¥”™¥áÑÕÉ”•¹•É…Ñ½È¥Ì¹½Ð‰åÑ”É•ÁÉ½‘Õ¥‰±”ˆ¤(€€€€€€€™¥áÑÕÉ•l‰‰åÑ•}É•ÁÉ½‘Õ¥‰±•}Í•½¹‘}•¹•É…Ñ¥½¸‰t€ôQÉÕ”(€€€€€€€É•½É‘Í}Á…Ñ €ôÉ½½Ð€¼€‰É•½É‘Ì¹©Í½¸ˆ(€€€€€€€ÁÉ•™±¥¡Ñ}Á…Ñ €ôÉ½½Ð€¼€‰ÁÉ•™±¥¡Ð¹©Í½¸ˆ(€€€€€€€•áÁ•Ñ•‘}Á…Ñ €ôÉ½½Ð€¼€‰•áÁ•Ñ•¹©Í½¸ˆ(€€€€€€€•áÁ•Ñ•‘}Á…Ñ ¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡ì‰¥¹¥Ñ¥…°ˆèì‰½É‘•É•‘}É•½É‘}µ…¹¥™•ÍÑ}¡…Í ˆè€‰Á•¹‘¥¹œ‰õô¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€€€€€ÁÉ•Á…É•}µ€ôl(€€€€€€€€€€€…ÉÌ¹ÁåÑ¡½¸°ÍÑÈ¡ÍÉ¥ÁÐ¤°€ˆ´µÝ½É­•Èˆ°€ˆ´µ½¡½ÉÐˆ°€‰ÁÉ•Á…É”ˆ°(€€€€€€€€€€€€ˆ´µÙ…É¥…¹Ðˆ°€‰…¹‘¥‘…Ñ”ˆ°€ˆ´µ½‘”µÉ½½Ðˆ°…ÉÌ¹…¹‘¥‘…Ñ•}É½½Ð°(€€€€€€€€€€€€ˆ´µ½ÕÑÁÕÐµÉ½½Ðˆ°ÍÑÈ¡É½½Ð€¼€‰ÁÉ•Á…É”µ‘ˆˆ¤°€ˆ´µ½ÉÁÕÌˆ°ÍÑÈ¡¥¹¥Ñ¥…°¤°(€€€€€€€€€€€€ˆ´µÕÁ‘…Ñ•µ½ÉÁÕÌˆ°ÍÑÈ¡ÕÁ‘…Ñ•¤°€ˆ´µµ½‘•°µ‘¥Èˆ°…ÉÌ¹µ½‘•±}‘¥È°(€€€€€€€€€€€€ˆ´µÉ•½É‘Ìˆ°ÍÑÈ¡É•½É‘Í}Á…Ñ ¤°€ˆ´µ•áÁ•Ñ•µµ…¹¥™•ÍÐˆ°ÍÑÈ¡•áÁ•Ñ•‘}Á…Ñ ¤°(€€€€€€€€€€€€ˆ´µ½ÕÑÁÕÐˆ°ÍÑÈ¡ÁÉ•™±¥¡Ñ}Á…Ñ ¤°(€€€€€€€t(€€€€€€€€ŒAÉ•Á…É”™¥ÉÍÐÝ¥Ñ¡½ÕÐÑ¡”•áÁ•Ñ•µ¡…Í ¡•¬°Ñ¡•¸™É••é”¥ÑÌµ…¹¥™•ÍÐ¸(€€€€€€€•áÁ•Ñ•‘}Á…Ñ ¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡ì‰¥¹¥Ñ¥…°ˆèì‰½É‘•É•‘}É•½É‘}µ…¹¥™•ÍÑ}¡…Í ˆè€ˆ‰õô¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€€€€€ÁÉ•Á…É•}µ¹…ÁÁ•¹ ˆ´µÁÉ•Á…É”µÕ¹¡•­•ˆ¤(€€€€€€€ÍÕ‰ÁÉ½•ÍÌ¹ÉÕ¸¡ÁÉ•Á…É•}µ°¡•¬õQÉÕ”°Ñ¥µ•½ÕÐôÐàÀ¤(€€€€€€€ÁÉ•™±¥¡Ð€ô©Í½¸¹±½…‘Ì¡ÁÉ•™±¥¡Ñ}Á…Ñ ¹É•…‘}Ñ•áÐ¡•¹½‘¥¹œô‰ÕÑ˜´àˆ¤¤(€€€€€€€•áÁ•Ñ•‘}Á…Ñ ¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡ÁÉ•™±¥¡Ð¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€€€€€¥˜ÁÉ•™±¥¡Ñl‰¥¹¥Ñ¥…°‰ul‰É•½É‘}½Õ¹Ð‰t€„ôQ=Q1}I=ILè(€€€€€€€€€€€É…¥Í”IÕ¹Ñ¥µ•ÉÉ½È ‰ÁÉ•™±¥¡Ð‘¥¹½ÐÁÉ½‘Õ”€ÐÀÀÉ•½É‘Ìˆ¤(€€€€€€€¥˜…ÉÌ¹ÁÉ•™±¥¡Ñ}½¹±äè(€€€€€€€€€€€½ÕÑÁÕÐ¹ÝÉ¥Ñ•}Ñ•áÐ (€€€€€€€€€€€€€€€©Í½¸¹‘ÕµÁÌ¡ì‰™¥áÑÕÉ”ˆè™¥áÑÕÉ”°€‰ÁÉ•™±¥¡ÐˆèÁÉ•™±¥¡Ñô°•¹ÍÕÉ•}…Í¥¤õ…±Í”°¥¹‘•¹ÐôÈ¤°(€€€€€€€€€€€€€€€•¹½‘¥¹œô‰ÕÑ˜´àˆ°(€€€€€€€€€€€€¤(€€€€€€€€€€€ÁÉ¥¹Ð¡©Í½¸¹‘ÕµÁÌ¡ì‰ÁÉ•™±¥¡ÐˆèÁÉ•™±¥¡Ñô°•¹ÍÕÉ•}…Í¥¤õ…±Í”¤¤(€€€€€€€€€€€É•ÑÕÉ¸€À((€€€€€€€…±±}É•ÍÕ±ÑÌè‘¥ÑmÍÑÈ°¹åt€ôì‰½É”ˆèíô°€‰”É”ˆèíõô(€€€€€€€ÑÉ…•}ÁÉ½‰•Ìè‘¥ÑmÍÑÈ°¹åt€ôíô(€€€€€€€™½È½¡½ÉÐ¥¸€ ‰½É”ˆ°€‰”É”ˆ¤è(€€€€€€€€€€€™½ÈÙ…É¥…¹Ð¥¸€ ‰‰…Í•±¥¹”ˆ°€‰…¹‘¥‘…Ñ”ˆ¤è(€€€€€€€€€€€€€€€½‘•}É½½Ð€ô…ÉÌ¹‰…Í•±¥¹•}É½½Ð¥˜Ù…É¥…¹Ð€ôô€‰‰…Í•±¥¹”ˆ•±Í”…ÉÌ¹…¹‘¥‘…Ñ•}É½½Ð(€€€€€€€€€€€€€€€Í…µÁ±•}É½½Ð€ôÉ½½Ð€¼€‰ÑÉ…”ˆ€¼½¡½ÉÐ€¼Ù…É¥…¹Ð(€€€€€€€€€€€€€€€½ÉÁÕÌ€ôÍ…µÁ±•}É½½Ð€¼€‰Í½ÕÉ”ˆ(€€€€€€€€€€€€€€€Í¡ÕÑ¥°¹½ÁåÑÉ•”¡¥¹¥Ñ¥…°°½ÉÁÕÌ¤(€€€€€€€€€€€€€€€Í…µÁ±•}½ÕÑÁÕÐ€ôÍ…µÁ±•}É½½Ð€¼€‰Í…µÁ±”¹©Í½¸ˆ(€€€€€€€€€€€€€€€µ€ôl(€€€€€€€€€€€€€€€€€€€…ÉÌ¹ÁåÑ¡½¸°ÍÑÈ¡ÍÉ¥ÁÐ¤°€ˆ´µÝ½É­•Èˆ°€ˆ´µÑÉ…”µÁÉ½‰”ˆ°(€€€€€€€€€€€€€€€€€€€€ˆ´µ½¡½ÉÐˆ°½¡½ÉÐ°€ˆ´µÙ…É¥…¹Ðˆ°Ù…É¥…¹Ð°(€€€€€€€€€€€€€€€€€€€€ˆ´µ½‘”µÉ½½Ðˆ°½‘•}É½½Ð°€ˆ´µ½ÕÑÁÕÐµÉ½½Ðˆ°ÍÑÈ¡Í…µÁ±•}É½½Ð€¼€‰‘ˆˆ¤°(€€€€€€€€€€€€€€€€€€€€ˆ´µ½ÉÁÕÌˆ°ÍÑÈ¡½ÉÁÕÌ¤°€ˆ´µÕÁ‘…Ñ•µ½ÉÁÕÌˆ°ÍÑÈ¡ÕÁ‘…Ñ•¤°(€€€€€€€€€€€€€€€€€€€€ˆ´µµ½‘•°µ‘¥Èˆ°…ÉÌ¹µ½‘•±}‘¥È°€ˆ´µÉ•½É‘Ìˆ°ÍÑÈ¡É•½É‘Í}Á…Ñ ¤°(€€€€€€€€€€€€€€€€€€€€ˆ´µ•áÁ•Ñ•µµ…¹¥™•ÍÐˆ°ÍÑÈ¡•áÁ•Ñ•‘}Á…Ñ ¤°€ˆ´µ½ÕÑÁÕÐˆ°ÍÑÈ¡Í…µÁ±•}½ÕÑÁÕÐ¤°(€€€€€€€€€€€€€€€t(€€€€€€€€€€€€€€€½µÁ±•Ñ•€ôÍÕ‰ÁÉ½•ÍÌ¹ÉÕ¸¡µ°Ñ¥µ•½ÕÐôÐàÀ°Ñ•áÐõQÉÕ”°…ÁÑÕÉ•}½ÕÑÁÕÐõQÉÕ”¤(€€€€€€€€€€€€€€€¥˜½µÁ±•Ñ•¹É•ÑÕÉ¹½‘”è(€€€€€€€€€€€€€€€€€€€É…¥Í”IÕ¹Ñ¥µ•ÉÉ½È¡˜‰ÑÉ…”Ý½É­•È™…¥±•èíµ‘õq¹í½µÁ±•Ñ•¹ÍÑ‘½ÕÑõq¹í½µÁ±•Ñ•¹ÍÑ‘•ÉÉôˆ¤(€€€€€€€€€€€€€€€Í…µÁ±”€ô©Í½¸¹±½…‘Ì¡Í…µÁ±•}½ÕÑÁÕÐ¹É•…‘}Ñ•áÐ¡•¹½‘¥¹œô‰ÕÑ˜´àˆ¤¤(€€€€€€€€€€€€€€€ÑÉ…•}ÁÉ½‰•Ím˜‰í½¡½ÉÑõ}íÙ…É¥…¹Ñô‰t€ôì(€€€€€€€€€€€€€€€€€€€€‰ÍÅ±}½Õ¹ÐˆèÍ…µÁ±•l‰ÍÅ±}½Õ¹Ð‰t°(€€€€€€€€€€€€€€€€€€€€‰Ñ½Ñ…±}¡…¹•ÌˆèÍ…µÁ±•l‰Ñ½Ñ…±}¡…¹•Ì‰t°(€€€€€€€€€€€€€€€ô(€€€€€€€É…Üè±¥ÍÑm‘¥ÑmÍÑÈ°¹åut€ômt(€€€€€€€Á…¥É}Á±…¸€ô€  ‰½É”ˆ°€Ä¤°€ ‰”É”ˆ°€Ä¤¤¥˜…ÉÌ¹Íµ½­”•±Í”€  ‰½É”ˆ°€Ô¤°€ ‰”É”ˆ°€Ì¤¤(€€€€€€€™½È½¡½ÉÐ°µ•…ÍÕÉ•‘}Á…¥ÉÌ¥¸Á…¥É}Á±…¸è(€€€€€€€€€€€Í…µÁ±•Ì€ôì‰‰…Í•±¥¹”ˆèmt°€‰…¹‘¥‘…Ñ”ˆèmuô(€€€€€€€€€€€™½ÈÁ…¥È¥¸É…¹”¡µ•…ÍÕÉ•‘}Á…¥ÉÌ€¬€Ä¤è(€€€€€€€€€€€€€€€½É‘•È€ô€ ‰‰…Í•±¥¹”ˆ°€‰…¹‘¥‘…Ñ”ˆ¤¥˜Á…¥È€”€È€ôô€À•±Í”€ ‰…¹‘¥‘…Ñ”ˆ°€‰‰…Í•±¥¹”ˆ¤(€€€€€€€€€€€€€€€Á…¥É}Ù…±Õ•Ìè‘¥ÑmÍÑÈ°‘¥ÑmÍÑÈ°¹åut€ôíô(€€€€€€€€€€€€€€€™½ÈÙ…É¥…¹Ð¥¸½É‘•Èè(€€€€€€€€€€€€€€€€€€€½‘•}É½½Ð€ô…ÉÌ¹‰…Í•±¥¹•}É½½Ð¥˜Ù…É¥…¹Ð€ôô€‰‰…Í•±¥¹”ˆ•±Í”…ÉÌ¹…¹‘¥‘…Ñ•}É½½Ð(€€€€€€€€€€€€€€€€€€€Í…µÁ±•}É½½Ð€ôÉ½½Ð€¼€‰ÉÕ¹Ìˆ€¼½¡½ÉÐ€¼˜‰íÁ…¥ÉôµíÙ…É¥…¹Ñôˆ(€€€€€€€€€€€€€€€€€€€½ÉÁÕÌ€ôÍ…µÁ±•}É½½Ð€¼€‰Í½ÕÉ”ˆ(€€€€€€€€€€€€€€€€€€€Í¡ÕÑ¥°¹½ÁåÑÉ•”¡¥¹¥Ñ¥…°°½ÉÁÕÌ¤(€€€€€€€€€€€€€€€€€€€Í…µÁ±•}½ÕÑÁÕÐ€ôÍ…µÁ±•}É½½Ð€¼€‰Í…µÁ±”¹©Í½¸ˆ(€€€€€€€€€€€€€€€€€€€µ€ôl(€€€€€€€€€€€€€€€€€€€€€€€…ÉÌ¹ÁåÑ¡½¸°ÍÑÈ¡ÍÉ¥ÁÐ¤°€ˆ´µÝ½É­•Èˆ°€ˆ´µ½¡½ÉÐˆ°½¡½ÉÐ°(€€€€€€€€€€€€€€€€€€€€€€€€ˆ´µÙ…É¥…¹Ðˆ°Ù…É¥…¹Ð°€ˆ´µ½‘”µÉ½½Ðˆ°½‘•}É½½Ð°(€€€€€€€€€€€€€€€€€€€€€€€€ˆ´µ½ÕÑÁÕÐµÉ½½Ðˆ°ÍÑÈ¡Í…µÁ±•}É½½Ð€¼€‰‘ˆˆ¤°€ˆ´µ½ÉÁÕÌˆ°ÍÑÈ¡½ÉÁÕÌ¤°(€€€€€€€€€€€€€€€€€€€€€€€€ˆ´µÕÁ‘…Ñ•µ½ÉÁÕÌˆ°ÍÑÈ¡ÕÁ‘…Ñ•¤°€ˆ´µµ½‘•°µ‘¥Èˆ°…ÉÌ¹µ½‘•±}‘¥È°(€€€€€€€€€€€€€€€€€€€€€€€€ˆ´µÉ•½É‘Ìˆ°ÍÑÈ¡É•½É‘Í}Á…Ñ ¤°€ˆ´µ•áÁ•Ñ•µµ…¹¥™•ÍÐˆ°ÍÑÈ¡•áÁ•Ñ•‘}Á…Ñ ¤°(€€€€€€€€€€€€€€€€€€€€€€€€ˆ´µ½ÕÑÁÕÐˆ°ÍÑÈ¡Í…µÁ±•}½ÕÑÁÕÐ¤°(€€€€€€€€€€€€€€€€€€€t(€€€€€€€€€€€€€€€€€€€½µÁ±•Ñ•€ôÍÕ‰ÁÉ½•ÍÌ¹ÉÕ¸¡µ°Ñ¥µ•½ÕÐôÐàÀ°Ñ•áÐõQÉÕ”°…ÁÑÕÉ•}½ÕÑÁÕÐõQÉÕ”¤(€€€€€€€€€€€€€€€€€€€¥˜½µÁ±•Ñ•¹É•ÑÕÉ¹½‘”è(€€€€€€€€€€€€€€€€€€€€€€€É…¥Í”IÕ¹Ñ¥µ•ÉÉ½È¡˜‰Ý½É­•È™…¥±•èíµ‘õq¹í½µÁ±•Ñ•¹ÍÑ‘½ÕÑõq¹í½µÁ±•Ñ•¹ÍÑ‘•ÉÉôˆ¤(€€€€€€€€€€€€€€€€€€€Í…µÁ±”€ô©Í½¸¹±½…‘Ì¡Í…µÁ±•}½ÕÑÁÕÐ¹É•…‘}Ñ•áÐ¡•¹½‘¥¹œô‰ÕÑ˜´àˆ¤¤(€€€€€€€€€€€€€€€€€€€Í…µÁ±”¹ÕÁ‘…Ñ”¡ì‰Á…¥ÈˆèÁ…¥È°€‰Ý…ÉµÕÀˆèÁ…¥È€ôô€Áô¤(€€€€€€€€€€€€€€€€€€€Á…¥É}Ù…±Õ•ÍmÙ…É¥…¹Ñt€ôÍ…µÁ±”(€€€€€€€€€€€€€€€¥˜Á…¥Èè(€€€€€€€€€€€€€€€€€€€™½ÈÙ…É¥…¹Ð¥¸€ ‰‰…Í•±¥¹”ˆ°€‰…¹‘¥‘…Ñ”ˆ¤è(€€€€€€€€€€€€€€€€€€€€€€€Í…µÁ±•ÍmÙ…É¥…¹Ñt¹…ÁÁ•¹¡Á…¥É}Ù…±Õ•ÍmÙ…É¥…¹Ñt¤(€€€€€€€€€€€€€€€€€€€€€€€É…Ü¹…ÁÁ•¹¡Á…¥É}Ù…±Õ•ÍmÙ…É¥…¹Ñt¤(€€€€€€€€€€€€€€€¥˜Ñ¥µ”¹µ½¹½Ñ½¹¥Œ ¤€´ÍÑ…ÉÑ•€ø€ÈÜÀÀè(€€€€€€€€€€€€€€€€€€€É…¥Í”Q¥µ•½ÕÑÉÉ½È ‰½µ‰¥¹•‰•¹¡µ…É¬•á••‘•€ÐÔµ¥¹ÕÑ•Ìˆ¤(€€€€€€€€€€€‰…Í•±¥¹”€ôÍ…µÁ±•Íl‰‰…Í•±¥¹”‰t(€€€€€€€€€€€…¹‘¥‘…Ñ”€ôÍ…µÁ±•Íl‰…¹‘¥‘…Ñ”‰t(€€€€€€€€€€€É•‘ÕÑ¥½¹Ì€ôì(€€€€€€€€€€€€€€€­•äè}É•‘ÕÑ¥½¸¡‰…Í•±¥¹”°…¹‘¥‘…Ñ”°­•ä¤(€€€€€€€€€€€€€€€™½È­•ä¥¸€ ‰½±‘}Ý…±±}Í•½¹‘Ìˆ°€‰½±‘}…Ñ…±½}Í•½¹‘Ìˆ°€‰ÕÁ‘…Ñ•}Ý…±±}Í•½¹‘Ìˆ°€‰ÕÁ‘…Ñ•}…Ñ…±½}Í•½¹‘Ìˆ¤(€€€€€€€€€€€ô(€€€€€€€€€€€…±±}É•ÍÕ±ÑÍm½¡½ÉÑt€ôì(€€€€€€€€€€€€€€€€‰‰…Í•±¥¹”ˆè}ÍÕµµ…Éä¡‰…Í•±¥¹”¤°(€€€€€€€€€€€€€€€€‰…¹‘¥‘…Ñ”ˆè}ÍÕµµ…Éä¡…¹‘¥‘…Ñ”¤°(€€€€€€€€€€€€€€€€‰Á…¥É•‘}µ•‘¥…¹}É•‘ÕÑ¥½¹ÌˆèÉ•‘ÕÑ¥½¹Ì°(€€€€€€€€€€€ô((€€€€€€€½É”€ô…±±}É•ÍÕ±ÑÍl‰½É”‰t(€€€€€€€”É”€ô…±±}É•ÍÕ±ÑÍl‰”É”‰t(€€€€€€€…Ñ•Ì€ôì(€€€€€€€€€€€€‰½É•}…Ñ…±½}…Ñ}±•…ÍÑ|ÈÁ}Á•É•¹Ðˆè½É•l‰Á…¥É•‘}µ•‘¥…¹}É•‘ÕÑ¥½¹Ì‰ul‰½±‘}…Ñ…±½}Í•½¹‘Ì‰t€øô€À¸ÈÀ°(€€€€€€€€€€€€‰½É•}™Õ±±}…‘‘}…Ñ}±•…ÍÑ|ÄÁ}Á•É•¹Ðˆè½É•l‰Á…¥É•‘}µ•‘¥…¹}É•‘ÕÑ¥½¹Ì‰ul‰½±‘}Ý…±±}Í•½¹‘Ì‰t€øô€À¸ÄÀ°(€€€€€€€€€€€€‰½™™¥•}”É•}¹½Ñ}Í±½Ý•Èˆè”É•l‰Á…¥É•‘}µ•‘¥…¹}É•‘ÕÑ¥½¹Ì‰ul‰½±‘}Ý…±±}Í•½¹‘Ì‰t€øô€À¸À°(€€€€€€€€€€€€‰½É•}ÕÁ‘…Ñ•}…Ñ}±•…ÍÑ|ÄÕ}Á•É•¹Ðˆè½É•l‰Á…¥É•‘}µ•‘¥…¹}É•‘ÕÑ¥½¹Ì‰ul‰ÕÁ‘…Ñ•}Ý…±±}Í•½¹‘Ì‰t€øô€À¸ÄÔ°(€€€€€€€€€€€€‰½±‘}ÀäÕ}Ý¥Ñ¡¥¹|ÄÁ}Á•É•¹Ðˆè½É•l‰…¹‘¥‘…Ñ”‰ul‰½±‘}Ý…±±}Í•½¹‘Í}ÀäÔ‰t€ðô½É•l‰‰…Í•±¥¹”‰ul‰½±‘}Ý…±±}Í•½¹‘Í}ÀäÔ‰t€¨€Ä¸ÄÀ°(€€€€€€€€€€€€‰ÕÁ‘…Ñ•}ÀäÕ}Ý¥Ñ¡¥¹|ÄÁ}Á•É•¹Ðˆè½É•l‰…¹‘¥‘…Ñ”‰ul‰ÕÁ‘…Ñ•}Ý…±±}Í•½¹‘Í}ÀäÔ‰t€ðô½É•l‰‰…Í•±¥¹”‰ul‰ÕÁ‘…Ñ•}Ý…±±}Í•½¹‘Í}ÀäÔ‰t€¨€Ä¸ÄÀ°(€€€€€€€€€€€€‰ÉÍÍ}Ý¥Ñ¡¥¹|ÄÕ}Á•É•¹Ðˆèµ…à (€€€€€€€€€€€€€€€½É•l‰…¹‘¥‘…Ñ”‰ul‰Á•…­}ÉÍÍ}‰åÑ•Í}ÀäÔ‰t€¼½É•l‰‰…Í•±¥¹”‰ul‰Á•…­}ÉÍÍ}‰åÑ•Í}ÀäÔ‰t°(€€€€€€€€€€€€€€€”É•l‰…¹‘¥‘…Ñ”‰ul‰Á•…­}ÉÍÍ}‰åÑ•Í}ÀäÔ‰t€¼”É•l‰‰…Í•±¥¹”‰ul‰Á•…­}ÉÍÍ}‰åÑ•Í}ÀäÔ‰t°(€€€€€€€€€€€€¤€ðô€Ä¸ÄÔ°(€€€€€€€€€€€€‰‘‰}‰åÑ•Í}Ý¥Ñ¡¥¹|Õ}Á•É•¹Ðˆèµ…à (€€€€€€€€€€€€€€€½É•l‰…¹‘¥‘…Ñ”‰ul‰‘‰}‰åÑ•Í}ÀäÔ‰t€¼½É•l‰‰…Í•±¥¹”‰ul‰‘‰}‰åÑ•Í}ÀäÔ‰t°(€€€€€€€€€€€€€€€”É•l‰…¹‘¥‘…Ñ”‰ul‰‘‰}‰åÑ•Í}ÀäÔ‰t€¼”É•l‰‰…Í•±¥¹”‰ul‰‘‰}‰åÑ•Í}ÀäÔ‰t°(€€€€€€€€€€€€¤€ðô€Ä¸ÀÔ°(€€€€€€€ô(€€€€€€€É•ÍÕ±Ð€ôì(€€€€€€€€€€€€‰Ñ…Í¬ˆè€‰1IHµAI´ÀÄÄµHÈˆ°(€€€€€€€€€€€€‰™¥áÑÕÉ”ˆè™¥áÑÕÉ”°(€€€€€€€€€€€€‰ÁÉ•™±¥¡ÐˆèÁÉ•™±¥¡Ð°(€€€€€€€€€€€€‰ÁåÑ¡½¸ˆè…ÉÌ¹ÁåÑ¡½¸°(€€€€€€€€€€€€‰µ½‘•±}‘¥Èˆè…ÉÌ¹µ½‘•±}‘¥È°(€€€€€€€€€€€€‰‰…Í•±¥¹•}É½½Ðˆè…ÉÌ¹‰…Í•±¥¹•}É½½Ð°(€€€€€€€€€€€€‰…¹‘¥‘…Ñ•}É½½Ðˆè…ÉÌ¹…¹‘¥‘…Ñ•}É½½Ð°(€€€€€€€€€€€€‰Ý…ÉµÕÁ}Á…¥ÉÌˆè€Ä°(€€€€€€€€€€€€‰½É•}µ•…ÍÕÉ•‘}Á…¥ÉÌˆèÁ…¥É}Á±…¹lÁulÅt°(€€€€€€€€€€€€‰½™™¥•}”É•}µ•…ÍÕÉ•‘}Á…¥ÉÌˆèÁ…¥É}Á±…¹lÅulÅt°(€€€€€€€€€€€€‰½¡½ÉÑÌˆè…±±}É•ÍÕ±ÑÌ°(€€€€€€€€€€€€‰…Ñ•Ìˆè…Ñ•Ì°(€€€€€€€€€€€€‰ÑÉ…•}ÁÉ½‰•Í}•á±Õ‘•‘}™É½µ}Ñ¥µ¥¹œˆèÑÉ…•}ÁÉ½‰•Ì°(€€€€€€€€€€€€‰™½Éµ…°ˆè¹½Ð…ÉÌ¹Íµ½­”°(€€€€€€€€€€€€‰Á…ÍÌˆè…±°¡…Ñ•Ì¹Ù…±Õ•Ì ¤¤°(€€€€€€€€€€€€‰É…Ý}Í…µÁ±•ÌˆèÉ…Ü°(€€€€€€€ô(€€€€€€€½ÕÑÁÕÐ¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡É•ÍÕ±Ð°•¹ÍÕÉ•}…Í¥¤õ…±Í”°¥¹‘•¹ÐôÈ¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€€€€€ÁÉ¥¹Ð¡©Í½¸¹‘ÕµÁÌ¡ì‰Á…ÍÌˆèÉ•ÍÕ±Ñl‰Á…ÍÌ‰t°€‰…Ñ•Ìˆè…Ñ•Íô°•¹ÍÕÉ•}…Í¥¤õ…±Í”¤¤(€€€€€€€É•ÑÕÉ¸€À¥˜…ÉÌ¹Íµ½­”½ÈÉ•ÍÕ±Ñl‰Á…ÍÌ‰t•±Í”€Ä(()‘•˜µ…¥¸ ¤€´ø¥¹Ðè(€€€Á…ÉÍ•È€ô…ÉÁ…ÉÍ”¹ÉÕµ•¹ÑA…ÉÍ•È ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÝ½É­•Èˆ°…Ñ¥½¸ô‰ÍÑ½É•}ÑÉÕ”ˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÁÉ•Á…É”µÕ¹¡•­•ˆ°…Ñ¥½¸ô‰ÍÑ½É•}ÑÉÕ”ˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÑÉ…”µÁÉ½‰”ˆ°…Ñ¥½¸ô‰ÍÑ½É•}ÑÉÕ”ˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÁÉ•™±¥¡Ðµ½¹±äˆ°…Ñ¥½¸ô‰ÍÑ½É•}ÑÉÕ”ˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÍµ½­”ˆ°…Ñ¥½¸ô‰ÍÑ½É•}ÑÉÕ”ˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ½¡½ÉÐˆ°¡½¥•Ìô ‰ÁÉ•Á…É”ˆ°€‰½É”ˆ°€‰”É”ˆ¤¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÙ…É¥…¹Ðˆ°¡½¥•Ìô ‰‰…Í•±¥¹”ˆ°€‰…¹‘¥‘…Ñ”ˆ¤¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ½‘”µÉ½½Ðˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ½ÕÑÁÕÐµÉ½½Ðˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ½ÉÁÕÌˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÕÁ‘…Ñ•µ½ÉÁÕÌˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µµ½‘•°µ‘¥Èˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÉ•½É‘Ìˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ•áÁ•Ñ•µµ…¹¥™•ÍÐˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ½ÕÑÁÕÐˆ°É•ÅÕ¥É•õQÉÕ”¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µÁåÑ¡½¸ˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ‰…Í•±¥¹”µÉ½½Ðˆ¤(€€€Á…ÉÍ•È¹…‘‘}…ÉÕµ•¹Ð ˆ´µ…¹‘¥‘…Ñ”µÉ½½Ðˆ¤(€€€…ÉÌ€ôÁ…ÉÍ•È¹Á…ÉÍ•}…ÉÌ ¤(€€€¥˜…ÉÌ¹Ý½É­•Èè(€€€€€€€¥˜…ÉÌ¹ÁÉ•Á…É•}Õ¹¡•­•è(€€€€€€€€€€€€ŒÉ••é”Ñ¡”•¹•É…Ñ•™¥áÑÕÉ”‰•™½É”‰…Í•±¥¹”½…¹‘¥‘…Ñ”½µÁ…É¥Í½¹Ì¸(€€€€€€€€€€€}½¹™¥ÕÉ”¡A…Ñ ¡…ÉÌ¹½‘•}É½½Ð¤¹É•Í½±Ù” ¤°A…Ñ ¡…ÉÌ¹½ÕÑÁÕÑ}É½½Ð¤¹É•Í½±Ù” ¤°A…Ñ ¡…ÉÌ¹µ½‘•±}‘¥È¤¹É•Í½±Ù” ¤¤(€€€€€€€€€€€¥¹¥Ñ¥…°°¥¹¥Ñ¥…±}µ…¹¥™•ÍÐ€ô}½É‘•É•‘}É•½É‘Ì¡A…Ñ ¡…ÉÌ¹½ÉÁÕÌ¤¹É•Í½±Ù” ¤¤(€€€€€€€€€€€ÕÁ‘…Ñ•°ÕÁ‘…Ñ•‘}µ…¹¥™•ÍÐ€ô}½É‘•É•‘}É•½É‘Ì¡A…Ñ ¡…ÉÌ¹ÕÁ‘…Ñ•‘}½ÉÁÕÌ¤¹É•Í½±Ù” ¤¤(€€€€€€€€€€€¡…¹•‘}Á…Ñ¡Ì€ôì(€€€€€€€€€€€€€€€€¨¡˜‰Ý½É½…Ñ…±½œµÕ¥‘”µí¥¹‘•àèÀÉ‘ô¹‘½àˆ™½È¥¹‘•à¥¸UAQ}=`¤°(€€€€€€€€€€€€€€€€¨¡˜‰Á½Ý•ÉÁ½¥¹Ð½…Ñ…±½œµÁ±…¸µí¥¹‘•àèÀÉ‘ô¹ÁÁÑàˆ™½È¥¹‘•à¥¸UAQ}AAQ`¤°(€€€€€€€€€€€ô(€€€€€€€€€€€¡…¹•€ôm¥Ñ•´™½È¥Ñ•´¥¸ÕÁ‘…Ñ•¥˜¥Ñ•µl‰µ•Ñ…‘…Ñ„‰ul‰Á…Ñ ‰t¹ÍÁ±¥Ð ˆ¼ˆ°€Ä¥l´Åt¥¸¡…¹•‘}Á…Ñ¡Ít(€€€€€€€€€€€½±‘}¥‘Ì€ôm¥Ñ•µl‰¥‰t™½È¥Ñ•´¥¸¥¹¥Ñ¥…°¥˜¥Ñ•µl‰µ•Ñ…‘…Ñ„‰ul‰Á…Ñ ‰t¹ÍÁ±¥Ð ˆ¼ˆ°€Ä¥l´Åt¥¸¡…¹•‘}Á…Ñ¡Ít(€€€€€€€€€€€¥˜±•¸¡¡…¹•¤€„ôUAQ}I=IL½È±•¸¡½±‘}¥‘Ì¤€„ôUAQ}I=ILè(€€€€€€€€€€€€€€€É…¥Í”IÕ¹Ñ¥µ•ÉÉ½È¡˜‰ÕÁ‘…Ñ”É•½Éµ¥Íµ…Ñ èí±•¸¡¡…¹•¥ô½í±•¸¡½±‘}¥‘Ì¥ôˆ¤(€€€€€€€€€€€A…Ñ ¡…ÉÌ¹É•½É‘Ì¤¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡ì‰¥¹¥Ñ¥…°ˆè¥¹¥Ñ¥…°°€‰ÕÁ‘…Ñ•ˆè¡…¹•°€‰½±‘}¥‘Ìˆè½±‘}¥‘Íô°•¹ÍÕÉ•}…Í¥¤õ…±Í”¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€€€€€€€€€A…Ñ ¡…ÉÌ¹½ÕÑÁÕÐ¤¹ÝÉ¥Ñ•}Ñ•áÐ¡©Í½¸¹‘ÕµÁÌ¡ì‰¥¹¥Ñ¥…°ˆè¥¹¥Ñ¥…±}µ…¹¥™•ÍÐ°€‰ÕÁ‘…Ñ•ˆèÕÁ‘…Ñ•‘}µ…¹¥™•ÍÑô°•¹ÍÕÉ•}…Í¥¤õ…±Í”°¥¹‘•¹ÐôÈ¤°•¹½‘¥¹œô‰ÕÑ˜´àˆ¤(€€€€€€€€€€€É•ÑÕÉ¸€À(€€€€€€€É•ÑÕÉ¸}Ý½É­•È¡…ÉÌ¤(€€€É•ÅÕ¥É•€ô€¡…ÉÌ¹ÁåÑ¡½¸°…ÉÌ¹‰…Í•±¥¹•}É½½Ð°…ÉÌ¹…¹‘¥‘…Ñ•}É½½Ð°…ÉÌ¹µ½‘•±}‘¥È¤(€€€¥˜¹½Ð…±°¡É•ÅÕ¥É•¤è(€€€€€€€Á…ÉÍ•È¹•ÉÉ½È ˆ´µÁåÑ¡½¸°€´µ‰…Í•±¥¹”µÉ½½Ð°€´µ…¹‘¥‘…Ñ”µÉ½½Ð°…¹€´µµ½‘•°µ‘¥È…É”É•ÅÕ¥É•ˆ¤(€€€É•ÑÕÉ¸}ÉÕ¸¡…ÉÌ¤(()¥˜}}¹…µ•}|€ôô€‰}}µ…¥¹}|ˆè(€€€É…¥Í”MåÍÑ•µá¥Ð¡µ…¥¸ ¤¤(