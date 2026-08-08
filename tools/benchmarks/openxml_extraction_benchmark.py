from __future__ import annotations

import argparse
import json
import os
import statistics
import subprocess
import sys
import tempfile
import time
from pathlib import Path
from typing import Any


PACKAGE_ROOT = Path(".copilot/rag/gen_db/software_rag_tool")
FORMATS = ("docx", "pptx")
SIZES = {
    "small": (8, 3),
    "medium": (80, 16),
    "large": (320, 48),
}


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Compare origin/main legacy Office extraction with direct Open XML "
            "extraction in isolated Python processes."
        )
    )
    parser.add_argument("--legacy-root", type=Path)
    parser.add_argument("--poc-root", type=Path)
    parser.add_argument("--repeat", type=int, default=5)
    parser.add_argument("--copies-per-size", type=int, default=3)
    parser.add_argument("--json-output", type=Path)
    parser.add_argument("--worker", action="store_true", help=argparse.SUPPRESS)
    parser.add_argument("--repo-root", type=Path, help=argparse.SUPPRESS)
    parser.add_argument("--fixture-root", type=Path, help=argparse.SUPPRESS)
    parser.add_argument("--format", choices=FORMATS, help=argparse.SUPPRESS)
    parser.add_argument(
        "--operation",
        choices=("benchmark", "inspect"),
        help=argparse.SUPPRESS,
    )
    args = parser.parse_args()

    if args.worker:
        print(
            json.dumps(
                _run_worker(
                    repo_root=_required(args.repo_root, "--repo-root"),
                    fixture_root=_required(args.fixture_root, "--fixture-root"),
                    file_format=_required(args.format, "--format"),
                    operation=_required(args.operation, "--operation"),
                    repeat=args.repeat,
                ),
                ensure_ascii=False,
            )
        )
        return 0

    legacy_root = _validated_repo_root(
        _required(args.legacy_root, "--legacy-root")
    )
    poc_root = _validated_repo_root(_required(args.poc_root, "--poc-root"))
    if args.repeat < 1:
        parser.error("--repeat must be at least 1")
    if args.copies_per_size < 1:
        parser.error("--copies-per-size must be at least 1")

    with tempfile.TemporaryDirectory(prefix="lrr-perf-006-") as temporary:
        fixture_root = Path(temporary)
        _generate_fixtures(fixture_root, args.copies_per_size)
        result = {
            "python": sys.executable,
            "repeat": args.repeat,
            "copies_per_size": args.copies_per_size,
            "benchmark": {},
            "comparison": {},
        }
        for file_format in FORMATS:
            legacy_benchmark = _invoke_worker(
                repo_root=legacy_root,
                fixture_root=fixture_root,
                file_format=file_format,
                operation="benchmark",
                repeat=args.repeat,
            )
            poc_benchmark = _invoke_worker(
                repo_root=poc_root,
                fixture_root=fixture_root,
                file_format=file_format,
                operation="benchmark",
                repeat=args.repeat,
            )
            legacy_inspection = _invoke_worker(
                repo_root=legacy_root,
                fixture_root=fixture_root,
                file_format=file_format,
                operation="inspect",
                repeat=args.repeat,
            )
            poc_inspection = _invoke_worker(
                repo_root=poc_root,
                fixture_root=fixture_root,
                file_format=file_format,
                operation="inspect",
                repeat=args.repeat,
            )
            result["benchmark"][file_format] = _benchmark_comparison(
                legacy_benchmark,
                poc_benchmark,
            )
            result["comparison"][file_format] = _compare_documents(
                legacy_inspection["documents"],
                poc_inspection["documents"],
            )

    if args.json_output:
        args.json_output.parent.mkdir(parents=True, exist_ok=True)
        args.json_output.write_text(
            json.dumps(result, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    _print_summary(result)
    return 0


def _run_worker(
    *,
    repo_root: Path,
    fixture_root: Path,
    file_format: str,
    operation: str,
    repeat: int,
) -> dict[str, Any]:
    sys.path.insert(0, str(repo_root.resolve() / PACKAGE_ROOT))
    from software_rag_tool.extractors import extract_sections

    paths = sorted(fixture_root.glob(f"*.{file_format}"))
    if not paths:
        raise RuntimeError(f"no .{file_format} fixtures found in {fixture_root}")

    def extract(path: Path) -> list[Any]:
        return extract_sections(
            path,
            chunk_max_chars=10_000_000,
            chunk_overlap=0,
        )

    if operation == "inspect":
        documents: dict[str, dict[str, Any]] = {}
        for path in paths:
            sections = extract(path)
            documents[path.name] = {
                "titles": [section.title for section in sections],
                "texts": [section.text for section in sections],
                "characters": sum(len(section.text) for section in sections),
            }
        return {"documents": documents}

    for path in paths:
        extract(path)
    samples: list[float] = []
    character_count = 0
    for _ in range(repeat):
        started = time.perf_counter()
        character_count = 0
        for path in paths:
            sections = extract(path)
            character_count += sum(len(section.text) for section in sections)
        samples.append(time.perf_counter() - started)
    median_seconds = statistics.median(samples)
    return {
        "documents": len(paths),
        "total_bytes": sum(path.stat().st_size for path in paths),
        "characters_per_run": character_count,
        "samples_seconds": samples,
        "median_seconds": median_seconds,
        "throughput_files_per_second": len(paths) / median_seconds,
    }


def _invoke_worker(
    *,
    repo_root: Path,
    fixture_root: Path,
    file_format: str,
    operation: str,
    repeat: int,
) -> dict[str, Any]:
    environment = os.environ.copy()
    environment["PYTHONUTF8"] = "1"
    command = [
        sys.executable,
        str(Path(__file__).resolve()),
        "--worker",
        "--repo-root",
        str(repo_root),
        "--fixture-root",
        str(fixture_root),
        "--format",
        file_format,
        "--operation",
        operation,
        "--repeat",
        str(repeat),
    ]
    completed = subprocess.run(
        command,
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        encoding="utf-8",
        env=environment,
    )
    return json.loads(completed.stdout)


def _benchmark_comparison(
    legacy: dict[str, Any],
    poc: dict[str, Any],
) -> dict[str, Any]:
    legacy_seconds = float(legacy["median_seconds"])
    poc_seconds = float(poc["median_seconds"])
    return {
        "documents": legacy["documents"],
        "total_bytes": legacy["total_bytes"],
        "legacy": legacy,
        "poc": poc,
        "reduction_percent": (legacy_seconds - poc_seconds)
        / legacy_seconds
        * 100.0,
    }


def _compare_documents(
    legacy_documents: dict[str, dict[str, Any]],
    poc_documents: dict[str, dict[str, Any]],
) -> dict[str, Any]:
    comparisons: list[dict[str, Any]] = []
    marker_misses = 0
    for filename in sorted(legacy_documents):
        legacy = legacy_documents[filename]
        poc = poc_documents[filename]
        legacy_text = "\n".join(legacy["texts"])
        poc_text = "\n".join(poc["texts"])
        expected_markers = _expected_markers(filename)
        missing_markers = [
            marker for marker in expected_markers if marker not in poc_text
        ]
        marker_misses += len(missing_markers)
        legacy_lines = {line for line in legacy_text.splitlines() if line.strip()}
        poc_lines = {line for line in poc_text.splitlines() if line.strip()}
        comparisons.append(
            {
                "filename": filename,
                "legacy_characters": len(legacy_text),
                "poc_characters": len(poc_text),
                "normalized_match": legacy_text == poc_text,
                "legacy_section_count": len(legacy["texts"]),
                "poc_section_count": len(poc["texts"]),
                "legacy_titles": legacy["titles"],
                "poc_titles": poc["titles"],
                "legacy_only_lines": sorted(legacy_lines - poc_lines)[:8],
                "poc_only_lines": sorted(poc_lines - legacy_lines)[:8],
                "missing_markers": missing_markers,
            }
        )
    return {
        "documents": comparisons,
        "normalized_matches": sum(
            int(item["normalized_match"]) for item in comparisons
        ),
        "total_documents": len(comparisons),
        "required_marker_misses": marker_misses,
    }


def _generate_fixtures(root: Path, copies_per_size: int) -> None:
    from docx import Document
    from pptx import Presentation
    from pptx.util import Inches

    for size, (unit_count, secondary_count) in SIZES.items():
        for copy_index in range(copies_per_size):
            docx_marker = f"MARKER-DOCX-{size}-{copy_index}"
            docx_tail = f"TAIL-DOCX-{size}-{copy_index}"
            document = Document()
            for paragraph_index in range(unit_count):
                paragraph = document.add_paragraph()
                paragraph.add_run(
                    f"{docx_marker} paragraph {paragraph_index} 日本語 English 42 "
                )
                paragraph.add_run("split-run content")
            table = document.add_table(rows=secondary_count, cols=3)
            for row_index, row in enumerate(table.rows):
                for column_index, cell in enumerate(row.cells):
                    cell.text = (
                        f"table {row_index}:{column_index} {docx_marker}"
                    )
            document.add_paragraph(docx_tail)
            document.save(root / f"{size}-{copy_index}.docx")

            pptx_marker = f"MARKER-PPTX-{size}-{copy_index}"
            pptx_tail = f"TAIL-PPTX-{size}-{copy_index}"
            presentation = Presentation()
            for slide_index in range(secondary_count):
                slide = presentation.slides.add_slide(
                    presentation.slide_layouts[6]
                )
                box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.5), Inches(9), Inches(6)
                )
                text_frame = box.text_frame
                text_frame.clear()
                for paragraph_index in range(
                    max(1, unit_count // secondary_count)
                ):
                    paragraph = (
                        text_frame.paragraphs[0]
                        if paragraph_index == 0
                        else text_frame.add_paragraph()
                    )
                    paragraph.add_run().text = (
                        f"{pptx_marker} slide {slide_index} paragraph "
                        f"{paragraph_index} 日本語 English 42 "
                    )
                    paragraph.add_run().text = "split-run content"
            tail_box = presentation.slides[-1].shapes.add_textbox(
                Inches(0.5), Inches(6.5), Inches(9), Inches(0.5)
            )
            tail_box.text = pptx_tail
            presentation.save(root / f"{size}-{copy_index}.pptx")


def _expected_markers(filename: str) -> tuple[str, str]:
    size, copy_with_extension = filename.split("-", maxsplit=1)
    copy_index = copy_with_extension.split(".", maxsplit=1)[0]
    extension = Path(filename).suffix[1:].upper()
    return (
        f"MARKER-{extension}-{size}-{copy_index}",
        f"TAIL-{extension}-{size}-{copy_index}",
    )


def _validated_repo_root(path: Path) -> Path:
    resolved = path.resolve()
    extractor = resolved / PACKAGE_ROOT / "software_rag_tool/extractors.py"
    if not extractor.is_file():
        raise ValueError(f"repository root does not contain extractors.py: {resolved}")
    return resolved


def _required(value: Any, option: str) -> Any:
    if value is None:
        raise ValueError(f"{option} is required")
    return value


def _print_summary(result: dict[str, Any]) -> None:
    print("| Format | Docs | Bytes | Legacy median | PoC median | Reduction |")
    print("|---|---:|---:|---:|---:|---:|")
    for file_format in FORMATS:
        item = result["benchmark"][file_format]
        print(
            f"| {file_format.upper()} | {item['documents']} | "
            f"{item['total_bytes']} | "
            f"{item['legacy']['median_seconds']:.6f}s | "
            f"{item['poc']['median_seconds']:.6f}s | "
            f"{item['reduction_percent']:.2f}% |"
        )
    print()
    print("| Format | Exact normalized | Required marker misses |")
    print("|---|---:|---:|")
    for file_format in FORMATS:
        item = result["comparison"][file_format]
        print(
            f"| {file_format.upper()} | {item['normalized_matches']}/"
            f"{item['total_documents']} | {item['required_marker_misses']} |"
        )


if __name__ == "__main__":
    raise SystemExit(main())
