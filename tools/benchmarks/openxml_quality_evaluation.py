from __future__ import annotations

import argparse
import contextlib
import hashlib
import io
import json
import os
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree


PACKAGE_ROOT = Path(".copilot/rag/gen_db/software_rag_tool")
DOCUMENTS = (
    {
        "file": "Buzz Ocean Field Notes.docx",
        "questions": (
            ("Buzz Ocean Field Notes", "Buzz Ocean Field Notes"),
            (
                "What does Buzz Ocean Field Notes identify as the cause of Incident 2187?",
                "smuggler-deposited artifact",
            ),
        ),
    },
    {
        "file": "Guide to Non-Contact Observation.docx",
        "questions": (
            ("Guide to Non-Contact Observation", "Guide to Non-Contact Observation"),
            (
                "Which activities does Guide to Non-Contact Observation prohibit?",
                "resource extraction are prohibited",
            ),
        ),
    },
    {
        "file": "Incident 2187 Legal Review.docx",
        "questions": (
            ("Incident 2187 Legal Review", "Incident 2187 Legal Review"),
            (
                "What habitat does Incident 2187 Legal Review describe for Buzz?",
                "tidally heated subsurface ocean and steam valleys",
            ),
        ),
    },
    {
        "file": "Station Evacuation Plan.pptx",
        "questions": (
            ("Station Evacuation Plan", "Station Evacuation Plan"),
            (
                "Which species does Station Evacuation Plan say the Federation protects?",
                "protected species",
            ),
        ),
    },
    {
        "file": "Correction Workshop_ Names and Units.pptx",
        "questions": (
            ("Correction Workshop: Names and Units", "Correction Workshop: Names and Units"),
            (
                "What cloud-top gravity does Correction Workshop: Names and Units state?",
                "5.49 G, not 5.49 m/s²",
            ),
        ),
    },
    {
        "file": "Field Lexicon for BZ-01.pptx",
        "questions": (
            ("Field Lexicon for BZ-01", "Field Lexicon for BZ-01"),
            (
                "Where do the Dam Folk live according to Field Lexicon for BZ-01?",
                "live on Buzz, not on Fizzbuzz",
            ),
        ),
    },
)
WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
DRAWING_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
PRESENTATION_NS = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)


def main() -> int:
    parser = argparse.ArgumentParser(
        description=(
            "Compare legacy and direct Open XML extraction for readable text, "
            "chunk quality, and bounded retrieval on public representative files."
        )
    )
    parser.add_argument("--legacy-root", type=Path)
    parser.add_argument("--poc-root", type=Path)
    parser.add_argument("--corpus-root", type=Path)
    parser.add_argument("--model-dir", type=Path)
    parser.add_argument("--json-output", type=Path)
    parser.add_argument("--worker", action="store_true", help=argparse.SUPPRESS)
    parser.add_argument("--repo-root", type=Path, help=argparse.SUPPRESS)
    parser.add_argument("--db-root", type=Path, help=argparse.SUPPRESS)
    args = parser.parse_args()

    if args.worker:
        payload = _worker(
            _required(args.repo_root, "--repo-root"),
            _required(args.corpus_root, "--corpus-root"),
            _required(args.model_dir, "--model-dir"),
            _required(args.db_root, "--db-root"),
        )
        print(json.dumps(payload, ensure_ascii=False))
        return 0

    legacy_root = _repo_root(_required(args.legacy_root, "--legacy-root"))
    poc_root = _repo_root(_required(args.poc_root, "--poc-root"))
    corpus_root = _required(args.corpus_root, "--corpus-root").resolve()
    model_dir = _required(args.model_dir, "--model-dir").resolve()
    _resolve_documents(corpus_root)
    if not (model_dir / "model.onnx").is_file():
        raise ValueError(f"model directory is incomplete: {model_dir}")

    with tempfile.TemporaryDirectory(prefix="lrr-perf-006-quality-") as temporary:
        root = Path(temporary)
        legacy = _invoke_worker(
            legacy_root, corpus_root, model_dir, root / "legacy"
        )
        poc = _invoke_worker(poc_root, corpus_root, model_dir, root / "poc")
    result = _compare(legacy, poc)
    if args.json_output:
        args.json_output.parent.mkdir(parents=True, exist_ok=True)
        args.json_output.write_text(
            json.dumps(result, ensure_ascii=False, indent=2) + "\n",
            encoding="utf-8",
        )
    _print_summary(result)
    return 0 if result["gates"]["poc_only_new_misses"] == 0 else 1


def _worker(
    repo_root: Path,
    corpus_root: Path,
    model_dir: Path,
    db_root: Path,
) -> dict[str, Any]:
    repo_root = _repo_root(repo_root)
    os.environ["RAG_OUTPUT_ROOT"] = str(db_root.resolve())
    os.environ["RAG_MODELS_DIR"] = str(model_dir.resolve().parent)
    os.environ["RAG_DB_NAME"] = "openxml-quality"
    os.environ["CHROMA_COLLECTION"] = "openxml_quality"
    sys.path.insert(0, str(repo_root / PACKAGE_ROOT))

    from software_rag_tool import catalog, store
    from software_rag_tool.manifest import write_manifest
    from software_rag_tool.records import build_records_for_file

    paths = _resolve_documents(corpus_root)
    records: list[dict[str, Any]] = []
    documents: dict[str, dict[str, Any]] = {}
    for path in paths:
        file_records = build_records_for_file(
            corpus_root,
            path,
            source_id="openxml-quality",
            chunk_max_chars=1_400,
            chunk_overlap=160,
        )
        records.extend(file_records)
        documents[path.name] = {
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            "bytes": path.stat().st_size,
            "structure": _structure(path),
            "raw_sections": _raw_sections(path),
            "chunks": [record["text"] for record in file_records],
            "chunk_titles": [
                record["metadata"]["section_path"] for record in file_records
            ],
        }

    # Product APIs print batch progress. Keep worker stdout as one JSON value.
    with contextlib.redirect_stdout(io.StringIO()):
        store.upsert_records(records)
        catalog.upsert_records(records)
        write_manifest(len(records))
        queries = [_run_query(item, store, catalog) for item in _questions()]
    return {
        "repo_head": _git_head(repo_root),
        "record_count": len(records),
        "documents": documents,
        "queries": queries,
    }


def _run_query(
    item: dict[str, str],
    store: Any,
    catalog: Any,
) -> dict[str, Any]:
    question = item["question"]
    lanes = {
        "hybrid": store.query(question, top_k=8, explain=True),
        "exact": catalog.exact_search(question, top_k=8),
        "lexical": catalog.bm25_search(question, top_k=8),
        "dense": store.vector_query(question, top_k=8),
    }
    return {
        **item,
        "ranks": {
            lane: _must_hit_rank(rows, item["file"], item["marker"])
            for lane, rows in lanes.items()
        },
        "signals": {
            lane: _row_summaries(rows) for lane, rows in lanes.items()
        },
    }


def _raw_sections(path: Path) -> list[str]:
    if path.suffix.lower() == ".docx":
        try:
            from software_rag_tool.openxml_text import extract_docx_text
        except ModuleNotFoundError:
            from docx import Document

            document = Document(str(path))
            parts = [p.text for p in document.paragraphs if p.text.strip()]
            for table in document.tables:
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
            return ["\n".join(parts)]
        return [extract_docx_text(path)]

    try:
        from software_rag_tool.openxml_text import extract_pptx_slide_texts
    except ModuleNotFoundError:
        from pptx import Presentation

        presentation = Presentation(str(path))
        return [
            "\n".join(
                shape.text
                for shape in slide.shapes
                if hasattr(shape, "text") and shape.text.strip()
            )
            for slide in presentation.slides
        ]
    return extract_pptx_slide_texts(path)


def _structure(path: Path) -> dict[str, int]:
    with zipfile.ZipFile(path) as package:
        if path.suffix.lower() == ".docx":
            root = ElementTree.fromstring(package.read("word/document.xml"))
            return {
                "paragraphs": len(root.findall(f".//{{{WORD_NS}}}p")),
                "tables": len(root.findall(f".//{{{WORD_NS}}}tbl")),
            }
        slide_names = sorted(
            name
            for name in package.namelist()
            if name.startswith("ppt/slides/slide") and name.endswith(".xml")
        )
        paragraphs = 0
        shapes = 0
        for name in slide_names:
            root = ElementTree.fromstring(package.read(name))
            paragraphs += len(root.findall(f".//{{{DRAWING_NS}}}p"))
            shapes += len(root.findall(f".//{{{PRESENTATION_NS}}}sp"))
        return {
            "slides": len(slide_names),
            "paragraphs": paragraphs,
            "text_shapes": shapes,
        }


def _compare(legacy: dict[str, Any], poc: dict[str, Any]) -> dict[str, Any]:
    documents: list[dict[str, Any]] = []
    for definition in DOCUMENTS:
        filename = str(definition["file"])
        left = legacy["documents"][filename]
        right = poc["documents"][filename]
        raw_equal = [_semantic_text(v) for v in left["raw_sections"]] == [
            _semantic_text(v) for v in right["raw_sections"]
        ]
        chunks_equal = [_semantic_text(v) for v in left["chunks"]] == [
            _semantic_text(v) for v in right["chunks"]
        ]
        documents.append(
            {
                "file": filename,
                "sha256": right["sha256"],
                "bytes": right["bytes"],
                "structure": right["structure"],
                "legacy_raw_sections": len(left["raw_sections"]),
                "poc_raw_sections": len(right["raw_sections"]),
                "legacy_chunks": len(left["chunks"]),
                "poc_chunks": len(right["chunks"]),
                "naturalness": "equivalent" if raw_equal else "meaningful_regression",
                "chunk_quality": "equivalent" if chunks_equal else "review_required",
            }
        )

    comparisons: list[dict[str, Any]] = []
    new_misses = 0
    legacy_hits = 0
    poc_hits = 0
    for left, right in zip(legacy["queries"], poc["queries"]):
        legacy_hit = left["ranks"]["hybrid"] is not None
        poc_hit = right["ranks"]["hybrid"] is not None
        legacy_hits += int(legacy_hit)
        poc_hits += int(poc_hit)
        new_misses += int(legacy_hit and not poc_hit)
        comparisons.append(
            {
                "question": right["question"],
                "file": right["file"],
                "marker": right["marker"],
                "legacy_ranks": left["ranks"],
                "poc_ranks": right["ranks"],
                "legacy_signals": left["signals"],
                "poc_signals": right["signals"],
                "poc_only_new_miss": legacy_hit and not poc_hit,
            }
        )
    return {
        "settings": {
            "documents": len(DOCUMENTS),
            "questions": len(comparisons),
            "chunk_max_chars": 1_400,
            "chunk_overlap": 160,
            "top_k": 8,
            "same_model_tokenizer_and_modes": True,
            "temporary_databases_removed": True,
        },
        "legacy_head": legacy["repo_head"],
        "poc_head": poc["repo_head"],
        "documents": documents,
        "retrieval": comparisons,
        "gates": {
            "naturalness": (
                "PASS"
                if all(d["naturalness"] == "equivalent" for d in documents)
                else "FAIL"
            ),
            "chunk_quality": (
                "PASS"
                if all(d["chunk_quality"] == "equivalent" for d in documents)
                else "FAIL"
            ),
            "legacy_must_hits": legacy_hits,
            "poc_must_hits": poc_hits,
            "poc_only_new_misses": new_misses,
            "retrieval": "PASS" if new_misses == 0 else "FAIL",
        },
    }


def _questions() -> list[dict[str, str]]:
    return [
        {"file": str(item["file"]), "question": question, "marker": marker}
        for item in DOCUMENTS
        for question, marker in item["questions"]
    ]


def _must_hit_rank(
    rows: list[dict[str, Any]], target_file: str, marker: str
) -> int | None:
    for rank, row in enumerate(rows, start=1):
        metadata = row.get("metadata") or {}
        if Path(str(metadata.get("path") or "")).name != target_file:
            continue
        if marker.casefold() in str(row.get("text") or "").casefold():
            return rank
    return None


def _row_summaries(rows: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [
        {
            "rank": rank,
            "file": Path(str((row.get("metadata") or {}).get("path") or "")).name,
            "section": str((row.get("metadata") or {}).get("section_path") or ""),
            "signals": list(row.get("signals") or []),
        }
        for rank, row in enumerate(rows[:8], start=1)
    ]


def _semantic_text(value: str) -> str:
    return " ".join(value.split())


def _resolve_documents(root: Path) -> list[Path]:
    paths: list[Path] = []
    for definition in DOCUMENTS:
        matches = list(root.rglob(str(definition["file"])))
        if len(matches) != 1:
            raise ValueError(
                f"expected one public corpus file named {definition['file']!r}, "
                f"found {len(matches)}"
            )
        paths.append(matches[0])
    return paths


def _invoke_worker(
    repo_root: Path,
    corpus_root: Path,
    model_dir: Path,
    db_root: Path,
) -> dict[str, Any]:
    environment = os.environ.copy()
    environment["PYTHONUTF8"] = "1"
    completed = subprocess.run(
        [
            sys.executable,
            str(Path(__file__).resolve()),
            "--worker",
            "--repo-root",
            str(repo_root),
            "--corpus-root",
            str(corpus_root),
            "--model-dir",
            str(model_dir),
            "--db-root",
            str(db_root),
        ],
        check=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        text=True,
        encoding="utf-8",
        env=environment,
    )
    return json.loads(completed.stdout)


def _git_head(repo_root: Path) -> str:
    return subprocess.run(
        [
            "git",
            "-c",
            f"safe.directory={repo_root.as_posix()}",
            "rev-parse",
            "HEAD",
        ],
        cwd=repo_root,
        check=True,
        stdout=subprocess.PIPE,
        text=True,
        encoding="utf-8",
    ).stdout.strip()


def _repo_root(path: Path) -> Path:
    root = path.resolve()
    if not (root / PACKAGE_ROOT / "software_rag_tool/extractors.py").is_file():
        raise ValueError(f"repository root does not contain extractors.py: {root}")
    return root


def _required(value: Any, option: str) -> Any:
    if value is None:
        raise ValueError(f"{option} is required")
    return value


def _print_summary(result: dict[str, Any]) -> None:
    gates = result["gates"]
    print(
        f"Naturalness={gates['naturalness']} "
        f"Chunk={gates['chunk_quality']} Retrieval={gates['retrieval']}"
    )
    print(
        f"Must-hit legacy={gates['legacy_must_hits']}/{result['settings']['questions']} "
        f"poc={gates['poc_must_hits']}/{result['settings']['questions']} "
        f"poc-only-new-misses={gates['poc_only_new_misses']}"
    )


if __name__ == "__main__":
    raise SystemExit(main())
