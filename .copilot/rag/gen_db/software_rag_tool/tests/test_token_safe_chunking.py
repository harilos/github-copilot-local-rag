from __future__ import annotations

import json
import os
import sys
import tempfile
import types
import unittest
from pathlib import Path
from unittest import mock

from software_rag_tool import incremental
from software_rag_tool.chunking import chunk_text, normalize_text
from software_rag_tool.config import default_onnx_model_dir
from software_rag_tool.embeddings import (
    DocumentEmbeddingTokenLimitError,
    DocumentTokenBudget,
    SentenceTransformerEmbedder,
    _validate_document_inputs,
    get_document_token_budget,
)
from software_rag_tool.ingestion_paths import resolve_ingestion_scope
from software_rag_tool.records import (
    build_records_for_file,
    chunker_config,
    file_content_hash,
)


class CharacterTokenizer:
    """Small explicit test double; every Unicode code point is one token."""

    def __call__(
        self,
        text: str | list[str],
        *,
        add_special_tokens: bool = True,
        **_kwargs: object,
    ) -> dict[str, list[int] | list[list[int]]]:
        def encode(value: str) -> list[int]:
            tokens = list(range(len(value)))
            return [10_001, *tokens, 10_002] if add_special_tokens else tokens

        if isinstance(text, list):
            return {"input_ids": [encode(value) for value in text]}
        return {"input_ids": encode(text)}


def test_budget(
    *,
    target_tokens: int = 80,
    max_tokens: int = 96,
) -> DocumentTokenBudget:
    return DocumentTokenBudget(
        tokenizer=CharacterTokenizer(),
        document_prefix="doc: ",
        tokenizer_name="explicit-character-test-double",
        target_tokens=target_tokens,
        max_tokens=max_tokens,
    )


class TokenSafeChunkingTests(unittest.TestCase):
    def test_mixed_content_is_covered_and_deterministic(self) -> None:
        cases = {
            "japanese": "日本語の段落です。末尾識別語。" * 30,
            "english": "An English sentence carries a tail marker. " * 30,
            "mixed": "日本語 Alpha-42 English。\n次の行です。" * 25,
            "no-newline": "TOKENWITHOUTBREAK0123456789" * 35,
            "table": "|列A|列B|\n|---|---|\n|値A|値B|\n" * 25,
            "code": "```python\ndef value():\n return 'marker'\n```\n" * 20,
            "crlf": "first line\r\nsecond line\r\n\r\nthird line\r\n" * 25,
        }
        budget = test_budget()
        path = "資料/長い経路/sample.md"
        title = "Long heading 見出し"
        for label, original in cases.items():
            with self.subTest(label=label):
                expected = normalize_text(original)
                first = chunk_text(
                    title,
                    original,
                    max_chars=120,
                    overlap=30,
                    token_budget=budget,
                    embedding_path=path,
                    overlap_tokens=12,
                )
                second = chunk_text(
                    title,
                    original,
                    max_chars=120,
                    overlap=30,
                    token_budget=budget,
                    embedding_path=path,
                    overlap_tokens=12,
                )
                self.assertEqual(first, second)
                self.assertGreater(len(first), 1)
                self.assertEqual(len(expected), first[-1].source_end)
                self.assertTrue(
                    all(
                        budget.count_document(path, item.title, item.text)
                        <= budget.max_tokens
                        for item in first
                    )
                )
                covered = [False] * len(expected)
                for item in first:
                    for index in range(item.source_start, item.source_end):
                        covered[index] = True
                self.assertTrue(
                    all(
                        is_covered or expected[index].isspace()
                        for index, is_covered in enumerate(covered)
                    )
                )
                for previous, current in zip(first, first[1:]):
                    self.assertGreater(current.source_end, previous.source_end)
                    overlap_text = expected[
                        current.source_start : previous.source_end
                    ]
                    self.assertLessEqual(len(overlap_text), 30)
                    self.assertLessEqual(budget.count_body(overlap_text), 12)

    def test_boundary_preference_order_and_fallback(self) -> None:
        budget = test_budget(target_tokens=42, max_tokens=48)
        cases = (
            ("a" * 26 + "\n\n" + "b" * 60, "a" * 26),
            ("a" * 26 + ". " + "b" * 60, "a" * 26 + "."),
            ("a" * 26 + "\n" + "b" * 60, "a" * 26),
            ("あ" * 26 + "。" + "い" * 60, "あ" * 26 + "。"),
        )
        for text, expected_first in cases:
            with self.subTest(expected_first=repr(expected_first)):
                chunks = chunk_text(
                    "",
                    text,
                    max_chars=200,
                    overlap=0,
                    token_budget=budget,
                    embedding_path="",
                )
                self.assertEqual(expected_first, chunks[0].text)

        fallback = chunk_text(
            "",
            "x" * 100,
            max_chars=200,
            overlap=0,
            token_budget=budget,
            embedding_path="",
        )
        self.assertGreater(len(fallback), 1)
        self.assertNotIn("\n", fallback[0].text)
        self.assertLessEqual(
            budget.count_document("", fallback[0].title, fallback[0].text),
            budget.target_tokens,
        )

    def test_document_encoder_rejects_oversize_instead_of_truncating(self) -> None:
        tokenizer = CharacterTokenizer()
        _validate_document_inputs(tokenizer, ["short"], max_tokens=16)
        with self.assertRaisesRegex(
            DocumentEmbeddingTokenLimitError,
            "silent truncation is disabled",
        ):
            _validate_document_inputs(tokenizer, ["x" * 20], max_tokens=16)

    def test_sentence_transformer_keeps_query_model_limit_unchanged(self) -> None:
        class FakeVector:
            def tolist(self) -> list[float]:
                return [1.0]

        class FakeModel:
            def __init__(self) -> None:
                self.max_seq_length = 64
                self.tokenizer = CharacterTokenizer()

            def encode(self, texts: list[str], **_kwargs: object) -> list[FakeVector]:
                return [FakeVector() for _ in texts]

        model = FakeModel()
        module = types.SimpleNamespace(
            SentenceTransformer=lambda *_args, **_kwargs: model
        )
        shared_tokenizer = CharacterTokenizer()
        with (
            mock.patch.dict(sys.modules, {"sentence_transformers": module}),
            mock.patch(
                "software_rag_tool.embeddings.get_embedding_tokenizer",
                return_value=shared_tokenizer,
            ),
        ):
            embedder = SentenceTransformerEmbedder(
                "fixture-model",
                "document: ",
                "query: ",
            )
        self.assertEqual(64, model.max_seq_length)
        self.assertEqual(64, embedder.max_length)
        self.assertIs(shared_tokenizer, model.tokenizer)
        self.assertEqual([[1.0]], embedder.encode(["q" * 100], mode="query"))
        with self.assertRaises(DocumentEmbeddingTokenLimitError):
            embedder.encode(["d" * 100], mode="document")

    def test_real_ruri_tokenizer_enforces_384_for_edge_cases(self) -> None:
        model_dir = default_onnx_model_dir()
        if not model_dir.exists():
            self.skipTest(f"local Ruri tokenizer is unavailable: {model_dir}")
        budget = get_document_token_budget()
        cases = (
            "日本語の長文です。" * 500,
            "unbroken_identifier_0123456789" * 500,
            ("| a | b |\n|---|---|\n| value | marker |\n" * 200),
            ("```python\nprint('marker')\n```\n" * 200),
            ("日本語 English français 12345。\r\n" * 300),
        )
        path = "長いパス/" * 12 + "document.md"
        title = "長い見出し" * 10
        for text in cases:
            chunks = chunk_text(
                title,
                text,
                token_budget=budget,
                embedding_path=path,
            )
            self.assertTrue(chunks)
            self.assertLessEqual(
                max(
                    budget.count_document(path, item.title, item.text)
                    for item in chunks
                ),
                384,
            )


class IncrementalReplacementTests(unittest.TestCase):
    def test_failed_upgrade_keeps_old_fingerprint_and_retries(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary) / "input"
            root.mkdir()
            source = root / "document.txt"
            source.write_text("content", encoding="utf-8")
            scope = resolve_ingestion_scope(root)
            stored_path = scope.file(source).stored_path
            key = f"fixture:{stored_path}"
            old_config = {"max_chars": 1400, "overlap": 160}
            state = {
                "files": {
                    key: {
                        "content_hash": file_content_hash(source),
                        "chunker_config": old_config,
                        "record_ids": ["legacy-id"],
                        "status": "indexed",
                    }
                }
            }
            budget = test_budget()
            current_config = chunker_config(
                chunk_max_chars=1400,
                chunk_overlap=160,
                document_token_budget=budget,
            )
            with mock.patch.object(
                incremental,
                "build_records_for_file",
                side_effect=RuntimeError("injected extraction failure"),
            ):
                failed = incremental._prepare_file(
                    scope,
                    source,
                    "fixture",
                    state,
                    retry_errors=False,
                    document_token_budget=budget,
                    current_chunker_config=current_config,
                )
            self.assertEqual("error", failed["status"])
            incremental._record_error(state, failed)
            self.assertEqual(old_config, state["files"][key]["chunker_config"])
            self.assertEqual(
                current_config,
                state["files"][key]["failed_chunker_config"],
            )
            with mock.patch.object(
                incremental,
                "build_records_for_file",
                return_value=[{"id": "new-id"}],
            ) as build_records:
                retry = incremental._prepare_file(
                    scope,
                    source,
                    "fixture",
                    state,
                    retry_errors=False,
                    document_token_budget=budget,
                    current_chunker_config=current_config,
                )
            self.assertEqual("ready", retry["status"])
            self.assertEqual(["legacy-id"], retry["previous_record_ids"])
            build_records.assert_called_once()

    def test_old_chunker_records_are_replaced_then_second_add_skips(self) -> None:
        with tempfile.TemporaryDirectory() as temporary:
            workspace = Path(temporary)
            root = workspace / "input"
            root.mkdir()
            source = root / "long.txt"
            source.write_text("末尾識別語。" * 80, encoding="utf-8")
            output = workspace / "db"
            logs = output / "logs"
            logs.mkdir(parents=True)
            scope = resolve_ingestion_scope(root)
            stored_path = scope.file(source).stored_path
            state_key = f"fixture:{stored_path}"
            state = {
                "version": 2,
                "files": {
                    state_key: {
                        "source_id": "fixture",
                        "path": stored_path,
                        "stored_path": stored_path,
                        "resolved_root": str(scope.resolved_root),
                        "content_hash": file_content_hash(source),
                        "chunker_config": {
                            "max_chars": 1400,
                            "overlap": 160,
                        },
                        "record_ids": ["legacy-1", "legacy-2"],
                        "status": "indexed",
                    }
                },
                "ingestion": {},
            }
            (logs / "index_state.json").write_text(
                json.dumps(state),
                encoding="utf-8",
            )
            budget = test_budget()

            with (
                mock.patch.dict(
                    os.environ,
                    {"RAG_OUTPUT_ROOT": str(output)},
                    clear=False,
                ),
                mock.patch.object(
                    incremental,
                    "delete_ids",
                    side_effect=lambda values: len(values),
                ) as delete_ids,
                mock.patch.object(incremental, "delete_catalog_chunks"),
                mock.patch.object(
                    incremental,
                    "upsert_records",
                    side_effect=lambda values, **_kwargs: len(values),
                ) as upsert_records,
                mock.patch.object(incremental, "upsert_catalog_records"),
                mock.patch.object(incremental, "collection_count", return_value=1),
                mock.patch.object(incremental, "write_manifest"),
                mock.patch.object(
                    incremental,
                    "update_profile_from_clean",
                    return_value=False,
                ),
                mock.patch.object(incremental, "write_progress"),
                mock.patch.object(incremental, "emit_event"),
            ):
                first = incremental.add_or_update_root(
                    root=root,
                    source_id="fixture",
                    document_token_budget=budget,
                )
                self.assertEqual(1, first["indexed_files"])
                delete_ids.assert_called_once_with(["legacy-1", "legacy-2"])
                self.assertTrue(upsert_records.called)

                saved = json.loads(
                    (logs / "index_state.json").read_text(encoding="utf-8")
                )
                current_ids = saved["files"][state_key]["record_ids"]
                self.assertTrue(current_ids)
                self.assertNotIn("legacy-1", current_ids)
                self.assertNotIn("legacy-2", current_ids)
                self.assertEqual(len(current_ids), len(set(current_ids)))

                delete_ids.reset_mock()
                upsert_records.reset_mock()
                second = incremental.add_or_update_root(
                    root=root,
                    source_id="fixture",
                    document_token_budget=budget,
                )
                self.assertEqual(1, second["skipped_files"])
                delete_ids.assert_not_called()
                upsert_records.assert_not_called()

                changed_prefix_budget = DocumentTokenBudget(
                    tokenizer=budget.tokenizer,
                    document_prefix="changed document prefix: ",
                    tokenizer_name=budget.tokenizer_name,
                    target_tokens=budget.target_tokens,
                    max_tokens=budget.max_tokens,
                )
                third = incremental.add_or_update_root(
                    root=root,
                    source_id="fixture",
                    document_token_budget=changed_prefix_budget,
                )
                self.assertEqual(1, third["indexed_files"])
                delete_ids.assert_called_once_with(current_ids)
                self.assertTrue(upsert_records.called)


class ExtractorFormatTokenBudgetTests(unittest.TestCase):
    def test_representative_formats_keep_tail_and_fit_ruri_limit(self) -> None:
        model_dir = default_onnx_model_dir()
        if not model_dir.exists():
            self.skipTest(f"local Ruri tokenizer is unavailable: {model_dir}")
        try:
            from docx import Document
            from pptx import Presentation
            from pptx.util import Inches
            from pypdf import PdfWriter
            from pypdf.generic import (
                DecodedStreamObject,
                DictionaryObject,
                NameObject,
            )
        except ModuleNotFoundError as exc:
            self.skipTest(f"document fixture dependency unavailable: {exc}")

        budget = get_document_token_budget()
        with tempfile.TemporaryDirectory() as temporary:
            root = Path(temporary) / "fixtures"
            root.mkdir()
            common = (
                "STARTMARKER Structured content 日本語 English. " * 180
                + " ENDMARKER"
            )
            (root / "sample.md").write_text(
                "# Heading\n\n" + common,
                encoding="utf-8",
            )
            (root / "sample.txt").write_text(common, encoding="utf-8")

            doc = Document()
            doc.add_paragraph(common)
            doc.save(root / "sample.docx")

            presentation = Presentation()
            slide = presentation.slides.add_slide(
                presentation.slide_layouts[6]
            )
            box = slide.shapes.add_textbox(
                Inches(0.5),
                Inches(0.5),
                Inches(9),
                Inches(6),
            )
            box.text_frame.text = common
            presentation.save(root / "sample.pptx")

            writer = PdfWriter()
            page = writer.add_blank_page(width=612, height=792)
            font = DictionaryObject(
                {
                    NameObject("/Type"): NameObject("/Font"),
                    NameObject("/Subtype"): NameObject("/Type1"),
                    NameObject("/BaseFont"): NameObject("/Helvetica"),
                }
            )
            page[NameObject("/Resources")] = DictionaryObject(
                {
                    NameObject("/Font"): DictionaryObject(
                        {NameObject("/F1"): writer._add_object(font)}
                    )
                }
            )
            pdf_lines = [
                f"PDF STARTMARKER line {index} content ENDMARKER"
                for index in range(160)
            ]
            commands = ["BT /F1 8 Tf 36 756 Td 10 TL"]
            commands.extend(
                f"({line}) Tj T*" for line in pdf_lines
            )
            commands.append("ET")
            stream = DecodedStreamObject()
            stream.set_data("\n".join(commands).encode("ascii"))
            page[NameObject("/Contents")] = writer._add_object(stream)
            with (root / "sample.pdf").open("wb") as output:
                writer.write(output)

            for path in sorted(root.iterdir()):
                with self.subTest(extension=path.suffix):
                    records = build_records_for_file(
                        root,
                        path,
                        source_id="format-fixture",
                        document_token_budget=budget,
                    )
                    self.assertTrue(records)
                    self.assertIn(
                        "ENDMARKER",
                        "\n".join(str(item["text"]) for item in records),
                    )
                    self.assertTrue(
                        all(
                            int(item["metadata"]["embedding_token_count"])
                            <= 384
                            for item in records
                        )
                    )


if __name__ == "__main__":
    unittest.main()
