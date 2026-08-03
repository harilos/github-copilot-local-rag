from __future__ import annotations

import codecs
import platform
import shutil
import subprocess
import tempfile
import unicodedata
from pathlib import Path

from .chunking import DocumentTokenBudget, TextSection, chunk_text, normalize_text
from .structured_extraction import (
    ExtractionResult,
    StructureBlock,
    extract_document_structure,
)

SUPPORTED_EXTENSIONS = {
    ".md",
    ".txt",
    ".log",
    ".pdf",
    ".docx",
    ".doc",
    ".pptx",
    ".ppt",
    ".xlsx",
    ".py",
    ".js",
    ".jsx",
    ".ts",
    ".tsx",
    ".java",
    ".go",
    ".rs",
    ".cs",
    ".rb",
    ".php",
    ".sh",
    ".ps1",
    ".sql",
    ".json",
    ".yaml",
    ".yml",
    ".toml",
    ".ini",
}


class ConverterOutputDecodeError(RuntimeError):
    """Raised when converter output is not valid supported text."""


_CONVERTER_BOM_ENCODINGS = (
    (codecs.BOM_UTF32_LE, "utf-32"),
    (codecs.BOM_UTF32_BE, "utf-32"),
    (codecs.BOM_UTF8, "utf-8-sig"),
    (codecs.BOM_UTF16_LE, "utf-16"),
    (codecs.BOM_UTF16_BE, "utf-16"),
)


def extract_sections(
    path: Path,
    *,
    chunk_max_chars: int = 1400,
    chunk_overlap: int = 160,
    token_budget: DocumentTokenBudget | None = None,
    embedding_path: str = "",
) -> list[TextSection]:
    ext = path.suffix.lower()
    if ext in {".doc", ".ppt"}:
        return _extract_legacy_office(path, chunk_max_chars=chunk_max_chars, chunk_overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path)
    extraction = extract_document(path)
    return sections_from_extraction(
        extraction,
        chunk_max_chars=chunk_max_chars,
        chunk_overlap=chunk_overlap,
        token_budget=token_budget,
        embedding_path=embedding_path,
    )


def extract_document(
    path: Path,
    *,
    backend_policy: str = "legacy",
    docling_threads: int | None = None,
) -> ExtractionResult:
    if path.suffix.lower() in {".doc", ".ppt"}:
        try:
            text = ""
            if path.suffix.lower() == ".doc" and platform.system() == "Darwin":
                text = _convert_doc_with_textutil(path)
            if not text:
                text = _convert_with_libreoffice(path)
        except Exception as exc:
            return ExtractionResult(
                status="extract_error",
                backend="legacy-office-converter",
                backend_version="strict-text-v1",
                source_format=path.suffix.lower().lstrip("."),
                structure_origin="legacy",
                retryable=True,
                reason=f"legacy_office_{type(exc).__name__}",
            )
        if not text.strip():
            return ExtractionResult(
                status="zero_text",
                backend="legacy-office-converter",
                backend_version="strict-text-v1",
                source_format=path.suffix.lower().lstrip("."),
                structure_origin="legacy",
                retryable=True,
                reason="legacy_office_zero_text",
            )
        return ExtractionResult(
            status="indexed",
            backend="legacy-office-converter",
            backend_version="strict-text-v1",
            source_format=path.suffix.lower().lstrip("."),
            structure_origin="legacy",
            blocks=(
                StructureBlock(
                    title=path.name,
                    text=text,
                    kind="document",
                    structure_id="legacy-office-0001",
                    parent_section_id="legacy-office-0001",
                ),
            ),
        )
    return extract_document_structure(
        path,
        backend_policy=backend_policy,
        docling_threads=docling_threads,
    )


def sections_from_extraction(
    extraction: ExtractionResult,
    *,
    chunk_max_chars: int,
    chunk_overlap: int,
    token_budget: DocumentTokenBudget | None,
    embedding_path: str,
) -> list[TextSection]:
    if not extraction.is_indexed:
        return []
    sections: list[TextSection] = []
    for block in extraction.blocks:
        block_sections = chunk_text(
            block.title,
            block.text,
            max_chars=chunk_max_chars,
            overlap=chunk_overlap,
            token_budget=token_budget,
            embedding_path=embedding_path,
            preserve_layout=block.preserve_layout,
            section_metadata={
                "source_format": extraction.source_format,
                "structure_origin": extraction.structure_origin,
                "structure_kind": block.kind,
                "structure_id": block.structure_id,
                "parent_section_id": block.parent_section_id,
                "breadcrumb": block.breadcrumb,
                "page": block.page,
                "slide": block.slide,
                "sheet": block.sheet,
            },
        )
        sections.extend(
            _offset_sections(block_sections, source_offset=block.source_start)
        )
    return sections


def _offset_sections(
    sections: list[TextSection],
    *,
    source_offset: int,
) -> list[TextSection]:
    if not source_offset:
        return sections
    return [
        TextSection(
            title=section.title,
            text=section.text,
            source_start=source_offset + section.source_start,
            source_end=source_offset + section.source_end,
            source_format=section.source_format,
            structure_origin=section.structure_origin,
            structure_kind=section.structure_kind,
            structure_id=section.structure_id,
            parent_section_id=section.parent_section_id,
            breadcrumb=section.breadcrumb,
            page=section.page,
            slide=section.slide,
            sheet=section.sheet,
        )
        for section in sections
    ]


def _extract_plain(path: Path, *, chunk_max_chars: int, chunk_overlap: int, token_budget: DocumentTokenBudget | None = None, embedding_path: str = "") -> list[TextSection]:
    text = path.read_text(encoding="utf-8", errors="replace")
    return chunk_text(path.name, text, max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path)


def _extract_pdf(path: Path, *, chunk_max_chars: int, chunk_overlap: int, token_budget: DocumentTokenBudget | None = None, embedding_path: str = "") -> list[TextSection]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    sections: list[TextSection] = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        sections.extend(chunk_text(f"Page {i}", text, max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path))
    if sections:
        return sections
    return [TextSection(title=path.name, text="")]


def _extract_docx(path: Path, *, chunk_max_chars: int, chunk_overlap: int, token_budget: DocumentTokenBudget | None = None, embedding_path: str = "") -> list[TextSection]:
    from docx import Document

    doc = Document(str(path))
    parts: list[str] = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            parts.append(paragraph.text)
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return chunk_text(path.name, "\n".join(parts), max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path)


def _extract_pptx(path: Path, *, chunk_max_chars: int, chunk_overlap: int, token_budget: DocumentTokenBudget | None = None, embedding_path: str = "") -> list[TextSection]:
    from pptx import Presentation

    prs = Presentation(str(path))
    sections: list[TextSection] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
        sections.extend(chunk_text(f"Slide {i}", "\n".join(parts), max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path))
    return sections


def _extract_xlsx(path: Path, *, chunk_max_chars: int, chunk_overlap: int, token_budget: DocumentTokenBudget | None = None, embedding_path: str = "") -> list[TextSection]:
    from openpyxl import load_workbook

    wb = load_workbook(str(path), read_only=True, data_only=True)
    sections: list[TextSection] = []
    for sheet in wb.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(v) for v in row if v is not None and str(v).strip()]
            if values:
                rows.append("\t".join(values))
        sections.extend(chunk_text(sheet.title, "\n".join(rows), max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path))
    return sections


def _extract_legacy_office(path: Path, *, chunk_max_chars: int, chunk_overlap: int, token_budget: DocumentTokenBudget | None = None, embedding_path: str = "") -> list[TextSection]:
    if path.suffix.lower() == ".doc" and platform.system() == "Darwin":
        text = _convert_doc_with_textutil(path)
        if text:
            return chunk_text(path.name, text, max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path)

    text = _convert_with_libreoffice(path)
    if text:
        return chunk_text(path.name, text, max_chars=chunk_max_chars, overlap=chunk_overlap, token_budget=token_budget, embedding_path=embedding_path)
    raise RuntimeError(f"Failed to extract legacy Office file: {path}")


def _convert_doc_with_textutil(path: Path) -> str:
    exe = shutil.which("textutil")
    if not exe:
        return ""
    proc = subprocess.run(
        [exe, "-convert", "txt", "-stdout", str(path)],
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=60,
    )
    if proc.returncode != 0:
        return ""
    return normalize_text(_decode_converter_output(proc.stdout))


def _convert_with_libreoffice(path: Path) -> str:
    exe = _find_libreoffice()
    if not exe:
        return ""
    with tempfile.TemporaryDirectory(prefix="ac-rag-office-") as tmp:
        tmp_path = Path(tmp)
        profile = tmp_path / "profile"
        cmd = [
            exe,
            f"-env:UserInstallation={profile.resolve().as_uri()}",
            "--headless",
            "--convert-to",
            "txt:Text",
            "--outdir",
            str(tmp_path),
            str(path),
        ]
        proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, timeout=180)
        if proc.returncode != 0:
            return ""
        txt_files = sorted(tmp_path.glob("*.txt"))
        if not txt_files:
            return ""
        return normalize_text(_decode_converter_output(txt_files[0].read_bytes()))


def _decode_converter_output(data: bytes) -> str:
    for bom, encoding in _CONVERTER_BOM_ENCODINGS:
        if data.startswith(bom):
            try:
                text = data.decode(encoding, errors="strict")
            except UnicodeError as exc:
                raise ConverterOutputDecodeError(
                    f"converter output has an invalid {encoding} byte sequence"
                ) from exc
            return _validate_converter_text(text)

    try:
        text = data.decode("utf-8", errors="strict")
    except UnicodeDecodeError:
        pass
    else:
        return _validate_converter_text(text)

    try:
        text = data.decode("cp932", errors="strict")
    except UnicodeDecodeError as exc:
        raise ConverterOutputDecodeError(
            "converter output is neither valid UTF-8 nor valid CP932"
        ) from exc

    return _validate_converter_text(text)


def _validate_converter_text(text: str) -> str:
    if any(
        unicodedata.category(char) == "Cc" and char not in "\t\n\r"
        for char in text
    ):
        raise ConverterOutputDecodeError(
            "converter output contains non-text control characters"
        )
    return text


def _find_libreoffice() -> str | None:
    candidates = [
        shutil.which("soffice.exe"),
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return str(candidate)
    return None
