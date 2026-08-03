from __future__ import annotations

import codecs
import hashlib
import importlib.metadata
import importlib.util
import os
import re
import unicodedata
from dataclasses import dataclass, field, replace
from pathlib import Path
from typing import Any, Iterable


EXTRACTION_STATUS_SCHEMA_VERSION = "extraction-status-v1"
MARKDOWN_PARSER_VERSION = "native-markdown-structure-v3"
XLSX_PARSER_VERSION = "openpyxl-region-structure-v2"
PLAIN_TEXT_ENCODING_POLICY_VERSION = "bom-utf8-cp932-strict-v1"
DOCLING_OPTIONS_VERSION = (
    "docling-tree-cpu-no-ocr-explicit-pdf-artifacts-v3"
)
DOCLING_PIN = "2.117.0"
DOCLING_ARTIFACTS_ENV = "RAG_DOCLING_ARTIFACTS_PATH"


@dataclass(frozen=True)
class StructureBlock:
    title: str
    text: str
    kind: str = "text"
    structure_id: str = ""
    parent_section_id: str = ""
    breadcrumb: str = ""
    source_start: int = 0
    source_end: int = 0
    page: int | None = None
    slide: int | None = None
    sheet: str = ""
    preserve_layout: bool = False


@dataclass(frozen=True)
class ExtractionResult:
    status: str
    backend: str
    backend_version: str
    source_format: str
    structure_origin: str
    blocks: tuple[StructureBlock, ...] = ()
    retryable: bool = False
    reason: str = ""
    encoding: str = ""
    encoding_reason: str = ""
    replacement_count: int = 0
    fallback_reason: str = ""
    diagnostics: dict[str, Any] = field(default_factory=dict)

    @property
    def is_indexed(self) -> bool:
        return self.status == "indexed" and any(
            block.text.strip() for block in self.blocks
        )


class PlainTextDecodeError(RuntimeError):
    pass


_TEXT_BOM_ENCODINGS = (
    (codecs.BOM_UTF32_LE, "utf-32"),
    (codecs.BOM_UTF32_BE, "utf-32"),
    (codecs.BOM_UTF8, "utf-8-sig"),
    (codecs.BOM_UTF16_LE, "utf-16"),
    (codecs.BOM_UTF16_BE, "utf-16"),
)

_ATX_HEADING = re.compile(r"^[ ]{0,3}(#{1,6})(?:[ \t]+|$)(.*)$")
_SETEXT_HEADING = re.compile(r"^[ ]{0,3}(=+|-+)[ \t]*$")
_FENCE_START = re.compile(r"^[ ]{0,3}(`{3,}|~{3,})")
_LIST_LINE = re.compile(r"^[ \t]*(?:[-+*]|\d+[.)])[ \t]+")
_TABLE_DIVIDER = re.compile(
    r"^[ \t]*\|?[ \t]*:?-{3,}:?[ \t]*(?:\|[ \t]*:?-{3,}:?[ \t]*)+\|?[ \t]*$"
)


def package_version(name: str) -> str:
    try:
        return importlib.metadata.version(name)
    except importlib.metadata.PackageNotFoundError:
        return "unavailable"


def docling_available() -> bool:
    return importlib.util.find_spec("docling") is not None


def docling_pdf_artifacts_path() -> Path | None:
    raw = str(os.getenv(DOCLING_ARTIFACTS_ENV) or "").strip()
    if not raw:
        return None
    path = Path(raw).expanduser().resolve()
    return path if path.is_dir() else None


def docling_pdf_artifacts_identity() -> str:
    """Return a location-independent identity for local PDF model artifacts."""

    path = docling_pdf_artifacts_path()
    if path is None:
        return "unconfigured"
    digest = hashlib.sha256()
    try:
        files = sorted(value for value in path.rglob("*") if value.is_file())
        for value in files:
            stat = value.stat()
            digest.update(value.relative_to(path).as_posix().encode("utf-8"))
            digest.update(str(stat.st_size).encode("ascii"))
            with value.open("rb") as stream:
                for chunk in iter(lambda: stream.read(1024 * 1024), b""):
                    digest.update(chunk)
    except OSError:
        return "unreadable"
    return "sha256:" + digest.hexdigest()


def extract_document_structure(
    path: Path,
    *,
    backend_policy: str = "legacy",
    docling_threads: int | None = None,
) -> ExtractionResult:
    """Extract semantic blocks without importing query-side heavy packages."""

    if backend_policy not in {"auto", "docling", "legacy"}:
        raise ValueError(f"unsupported backend policy: {backend_policy}")
    ext = path.suffix.lower()
    source_format = ext.lstrip(".") or "unknown"
    if ext == ".md":
        if backend_policy == "legacy":
            plain = extract_plain_text(path)
            return replace(
                plain,
                backend="legacy-markdown-flat",
                source_format="md",
                structure_origin="legacy_flat",
            )
        return extract_native_markdown(path)
    if ext in {
        ".txt",
        ".log",
        ".py",
        ".js",
        ".jsx",
        ".ts",
        ".tsx",
        ".java",
        ".go",
        ".rs",
        ".cs",
        ".rb",
        ".php",
        ".sh",
        ".ps1",
        ".sql",
        ".json",
        ".yaml",
        ".yml",
        ".toml",
        ".ini",
    }:
        return extract_plain_text(path)
    if ext == ".xlsx":
        if backend_policy == "legacy":
            return extract_legacy_xlsx_structure(path)
        return extract_xlsx_structure(path)
    if ext in {".pdf", ".docx", ".pptx"}:
        pdf_artifacts_missing = (
            ext == ".pdf" and docling_pdf_artifacts_path() is None
        )
        if (
            backend_policy != "legacy"
            and docling_available()
            and not pdf_artifacts_missing
        ):
            result = extract_docling_tree(path, threads=docling_threads)
            if result.is_indexed or backend_policy == "docling":
                return result
            fallback_reason = result.reason or result.status
        elif backend_policy == "docling":
            return ExtractionResult(
                status="unsupported",
                backend="docling",
                backend_version="unavailable",
                source_format=source_format,
                structure_origin="docling_tree",
                retryable=False,
                reason=(
                    "docling_pdf_artifacts_unavailable_offline"
                    if pdf_artifacts_missing
                    else "docling_unavailable"
                ),
            )
        else:
            fallback_reason = (
                "legacy_requested"
                if backend_policy == "legacy"
                else (
                    "docling_pdf_artifacts_unavailable_offline"
                    if pdf_artifacts_missing
                    else "docling_unavailable"
                )
            )
        legacy = extract_legacy_office_structure(path)
        return ExtractionResult(
            **{
                **legacy.__dict__,
                "structure_origin": "legacy_fallback",
                "fallback_reason": fallback_reason,
            }
        )
    return ExtractionResult(
        status="unsupported",
        backend="none",
        backend_version="",
        source_format=source_format,
        structure_origin="unsupported",
        retryable=False,
        reason="unsupported_format",
    )


def extract_plain_text(path: Path) -> ExtractionResult:
    source_format = path.suffix.lower().lstrip(".") or "text"
    try:
        text, encoding, reason = decode_plain_text(path.read_bytes())
    except (OSError, PlainTextDecodeError) as exc:
        return ExtractionResult(
            status="extract_error",
            backend="plain-text",
            backend_version=PLAIN_TEXT_ENCODING_POLICY_VERSION,
            source_format=source_format,
            structure_origin="plain_text",
            retryable=True,
            reason=_error_code("plain_text", exc),
        )
    if not text.strip():
        return ExtractionResult(
            status="zero_text",
            backend="plain-text",
            backend_version=PLAIN_TEXT_ENCODING_POLICY_VERSION,
            source_format=source_format,
            structure_origin="plain_text",
            retryable=True,
            reason="empty_text",
            encoding=encoding,
            encoding_reason=reason,
        )
    block = StructureBlock(
        title=path.name,
        text=_normalize_line_endings(text),
        kind="text",
        structure_id="plain-0001",
        parent_section_id="plain-0001",
        source_start=0,
        source_end=len(text),
        preserve_layout=source_format in {
            "py",
            "js",
            "jsx",
            "ts",
            "tsx",
            "java",
            "go",
            "rs",
            "cs",
            "rb",
            "php",
            "sh",
            "ps1",
            "sql",
            "json",
            "yaml",
            "yml",
            "toml",
            "ini",
        },
    )
    return ExtractionResult(
        status="indexed",
        backend="plain-text",
        backend_version=PLAIN_TEXT_ENCODING_POLICY_VERSION,
        source_format=source_format,
        structure_origin="plain_text",
        blocks=(block,),
        encoding=encoding,
        encoding_reason=reason,
        replacement_count=0,
    )


def decode_plain_text(data: bytes) -> tuple[str, str, str]:
    for bom, encoding in _TEXT_BOM_ENCODINGS:
        if data.startswith(bom):
            try:
                text = data.decode(encoding, errors="strict")
            except UnicodeError as exc:
                raise PlainTextDecodeError(
                    f"invalid BOM-declared encoding: {encoding}"
                ) from exc
            return _validate_plain_text(text), encoding, "bom"
    try:
        text = data.decode("utf-8", errors="strict")
    except UnicodeDecodeError:
        pass
    else:
        return _validate_plain_text(text), "utf-8", "strict_utf8"
    try:
        text = data.decode("cp932", errors="strict")
    except UnicodeDecodeError as exc:
        raise PlainTextDecodeError(
            "input is neither strict UTF-8 nor strict CP932"
        ) from exc
    return _validate_plain_text(text), "cp932", "deterministic_fallback"


def _validate_plain_text(text: str) -> str:
    if "\ufffd" in text:
        raise PlainTextDecodeError("replacement characters are forbidden")
    if any(
        unicodedata.category(char) == "Cc" and char not in "\t\n\r\f"
        for char in text
    ):
        raise PlainTextDecodeError("input contains non-text control characters")
    return text


def extract_native_markdown(path: Path) -> ExtractionResult:
    try:
        text, encoding, reason = decode_plain_text(path.read_bytes())
    except (OSError, PlainTextDecodeError) as exc:
        return ExtractionResult(
            status="extract_error",
            backend="native-markdown",
            backend_version=MARKDOWN_PARSER_VERSION,
            source_format="md",
            structure_origin="native_markdown",
            retryable=True,
            reason=_error_code("markdown", exc),
        )
    text = _normalize_line_endings(text)
    blocks = tuple(_parse_markdown_blocks(path.name, text))
    if not any(block.text.strip() for block in blocks):
        return ExtractionResult(
            status="zero_text",
            backend="native-markdown",
            backend_version=MARKDOWN_PARSER_VERSION,
            source_format="md",
            structure_origin="native_markdown",
            retryable=True,
            reason="empty_markdown_body",
            encoding=encoding,
            encoding_reason=reason,
        )
    return ExtractionResult(
        status="indexed",
        backend="native-markdown",
        backend_version=MARKDOWN_PARSER_VERSION,
        source_format="md",
        structure_origin="native_markdown",
        blocks=blocks,
        encoding=encoding,
        encoding_reason=reason,
    )


def _parse_markdown_blocks(name: str, text: str) -> Iterable[StructureBlock]:
    lines = text.splitlines(keepends=True)
    offsets: list[int] = []
    offset = 0
    for line in lines:
        offsets.append(offset)
        offset += len(line)

    headings: list[str] = []
    buffer: list[str] = []
    buffer_start = 0
    ordinal = 0
    fence_marker = ""
    index = 0

    def title() -> str:
        return " > ".join(headings) if headings else name

    def flush(end_offset: int) -> StructureBlock | None:
        nonlocal buffer, buffer_start, ordinal
        body = "".join(buffer).strip("\n")
        buffer = []
        if not body.strip():
            return None
        ordinal += 1
        # Small chunks keep the full breadcrumb as their title, while their
        # stable parent identifies the top-level semantic section.  This lets
        # an opt-in small-to-big pass recover sibling evidence without making
        # every Markdown file one undifferentiated parent.
        parent = _stable_id(
            "md-parent",
            headings[0] if headings else name,
        )
        return StructureBlock(
            title=title(),
            text=body,
            kind=_markdown_block_kind(body),
            structure_id=f"md-{ordinal:04d}",
            parent_section_id=parent,
            breadcrumb=" > ".join(headings),
            source_start=buffer_start,
            source_end=end_offset,
            preserve_layout=True,
        )

    if lines and lines[0].strip() == "---":
        closing = _yaml_front_matter_closing(lines)
        if closing is not None:
            ordinal += 1
            end = offsets[closing] + len(lines[closing])
            yield StructureBlock(
                title="YAML front matter",
                text="".join(lines[: closing + 1]).strip("\n"),
                kind="yaml",
                structure_id=f"md-{ordinal:04d}",
                parent_section_id="md-front-matter",
                source_start=0,
                source_end=end,
                preserve_layout=True,
            )
            index = closing + 1
            buffer_start = end

    while index < len(lines):
        line = lines[index]
        current_offset = offsets[index]
        if fence_marker:
            buffer.append(line)
            stripped = line.lstrip(" ")
            if re.fullmatch(
                rf"{re.escape(fence_marker[0])}"
                rf"{{{len(fence_marker)},}}[ \t]*(?:\n)?",
                stripped,
            ):
                fence_marker = ""
            index += 1
            continue

        fence = _FENCE_START.match(line)
        if fence:
            if not buffer:
                buffer_start = current_offset
            fence_marker = fence.group(1)
            buffer.append(line)
            index += 1
            continue

        atx = _ATX_HEADING.match(line.rstrip("\n"))
        setext = (
            index + 1 < len(lines)
            and line.strip()
            and _SETEXT_HEADING.match(lines[index + 1].rstrip("\n"))
        )
        if atx or setext:
            block = flush(current_offset)
            if block:
                yield block
            if atx:
                level = len(atx.group(1))
                heading = re.sub(
                    r"[ \t]+#+[ \t]*$", "", atx.group(2)
                ).strip()
                next_index = index + 1
            else:
                level = 1 if lines[index + 1].lstrip().startswith("=") else 2
                heading = line.strip()
                next_index = index + 2
            headings[level - 1 :] = [heading or f"Section {level}"]
            index = next_index
            buffer_start = offsets[index] if index < len(lines) else len(text)
            continue

        if not buffer:
            buffer_start = current_offset
        buffer.append(line)
        index += 1

    block = flush(len(text))
    if block:
        yield block


def _yaml_front_matter_closing(lines: list[str]) -> int | None:
    closing = next(
        (
            candidate
            for candidate in range(1, min(len(lines), 201))
            if lines[candidate].strip() in {"---", "..."}
        ),
        None,
    )
    if closing is None:
        return None
    try:
        import yaml

        value = yaml.safe_load("".join(lines[1:closing]))
    except Exception:
        return None
    return closing if isinstance(value, dict) else None


def _markdown_block_kind(text: str) -> str:
    lines = [line for line in text.splitlines() if line.strip()]
    if not lines:
        return "text"
    if _FENCE_START.match(lines[0]):
        return "code"
    if all(_LIST_LINE.match(line) for line in lines):
        return "list"
    if len(lines) >= 2 and _TABLE_DIVIDER.match(lines[1]):
        return "table"
    return "section"


def extract_xlsx_structure(path: Path) -> ExtractionResult:
    try:
        from openpyxl import load_workbook

        formula_book = load_workbook(str(path), read_only=True, data_only=False)
        value_book = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        return ExtractionResult(
            status="extract_error",
            backend="openpyxl",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="xlsx_structure",
            retryable=True,
            reason=_error_code("xlsx_open", exc),
        )
    blocks: list[StructureBlock] = []
    missing_formula_cache = 0
    try:
        for formula_sheet in formula_book.worksheets:
            value_sheet = value_book[formula_sheet.title]
            region: list[str] = []
            region_number = 0
            for formula_row, value_row in zip(
                formula_sheet.iter_rows(), value_sheet.iter_rows()
            ):
                values: list[str] = []
                for formula_cell, value_cell in zip(formula_row, value_row):
                    value = value_cell.value
                    if formula_cell.data_type == "f" and value is None:
                        missing_formula_cache += 1
                        value = formula_cell.value
                    values.append("" if value is None else str(value))
                while values and not values[-1].strip():
                    values.pop()
                if not any(value.strip() for value in values):
                    if region:
                        region_number += 1
                        blocks.append(
                            _xlsx_region_block(
                                formula_sheet.title,
                                region_number,
                                region,
                            )
                        )
                        region = []
                    continue
                region.append("\t".join(values))
            if region:
                region_number += 1
                blocks.append(
                    _xlsx_region_block(
                        formula_sheet.title,
                        region_number,
                        region,
                    )
                )
    except Exception as exc:
        return ExtractionResult(
            status="extract_error",
            backend="openpyxl",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="xlsx_structure",
            retryable=True,
            reason=_error_code("xlsx_iterate", exc),
        )
    finally:
        formula_book.close()
        value_book.close()
    blocks = list(_with_canonical_offsets(blocks))
    diagnostics = {
        "formula_cache_missing": missing_formula_cache,
        "source_range_basis": "canonical_extracted_text",
    }
    if missing_formula_cache:
        return ExtractionResult(
            status="extract_error",
            backend="openpyxl",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="xlsx_structure",
            blocks=tuple(blocks),
            retryable=True,
            reason="formula_cache_missing",
            diagnostics=diagnostics,
        )
    if not blocks:
        return ExtractionResult(
            status="zero_text",
            backend="openpyxl",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="xlsx_structure",
            retryable=True,
            reason="empty_workbook",
            diagnostics=diagnostics,
        )
    return ExtractionResult(
        status="indexed",
        backend="openpyxl",
        backend_version=package_version("openpyxl"),
        source_format="xlsx",
        structure_origin="xlsx_structure",
        blocks=tuple(blocks),
        diagnostics=diagnostics,
    )


def extract_legacy_xlsx_structure(path: Path) -> ExtractionResult:
    """Preserve the pre-Phase-2 flat, data-only XLSX extraction contract."""

    try:
        from openpyxl import load_workbook

        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        return ExtractionResult(
            status="extract_error",
            backend="openpyxl-legacy",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="legacy_flat",
            retryable=True,
            reason=_error_code("xlsx_legacy_open", exc),
        )
    blocks: list[StructureBlock] = []
    try:
        for index, sheet in enumerate(workbook.worksheets, start=1):
            rows: list[str] = []
            for row in sheet.iter_rows(values_only=True):
                values = [
                    str(value)
                    for value in row
                    if value is not None and str(value).strip()
                ]
                if values:
                    rows.append("\t".join(values))
            text = "\n".join(rows)
            if text.strip():
                blocks.append(
                    StructureBlock(
                        title=sheet.title,
                        text=text,
                        kind="sheet",
                        structure_id=f"legacy-sheet-{index:04d}",
                        parent_section_id=f"legacy-sheet-{index:04d}",
                        sheet=sheet.title,
                        preserve_layout=True,
                    )
                )
    except Exception as exc:
        return ExtractionResult(
            status="extract_error",
            backend="openpyxl-legacy",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="legacy_flat",
            retryable=True,
            reason=_error_code("xlsx_legacy_iterate", exc),
        )
    finally:
        workbook.close()
    blocks = list(_with_canonical_offsets(blocks))
    if not blocks:
        return ExtractionResult(
            status="zero_text",
            backend="openpyxl-legacy",
            backend_version=package_version("openpyxl"),
            source_format="xlsx",
            structure_origin="legacy_flat",
            retryable=True,
            reason="empty_workbook",
        )
    return ExtractionResult(
        status="indexed",
        backend="openpyxl-legacy",
        backend_version=package_version("openpyxl"),
        source_format="xlsx",
        structure_origin="legacy_flat",
        blocks=tuple(blocks),
        diagnostics={"source_range_basis": "canonical_extracted_text"},
    )


def _xlsx_region_block(
    sheet: str,
    region_number: int,
    rows: list[str],
) -> StructureBlock:
    parent = _stable_id("xlsx", f"{sheet}:{region_number}")
    return StructureBlock(
        title=f"{sheet} > Table {region_number}",
        text="\n".join(rows),
        kind="table",
        structure_id=parent,
        parent_section_id=parent,
        breadcrumb=sheet,
        sheet=sheet,
        preserve_layout=True,
    )


_DOCLING_CONVERTER: Any | None = None


def extract_docling_tree(
    path: Path,
    *,
    threads: int | None = None,
) -> ExtractionResult:
    source_format = path.suffix.lower().lstrip(".")
    try:
        converter = _get_docling_converter(threads=threads)
        conversion = converter.convert(path, raises_on_error=False)
        document = getattr(conversion, "document", None)
        status = str(getattr(conversion, "status", "")).lower()
        if document is None or (status and "failure" in status):
            return ExtractionResult(
                status="extract_error",
                backend="docling",
                backend_version=package_version("docling"),
                source_format=source_format,
                structure_origin="docling_tree",
                retryable=True,
                reason="docling_conversion_failed",
            )
        blocks = tuple(
            _with_canonical_offsets(
                _docling_blocks(document, source_format)
            )
        )
    except Exception as exc:
        return ExtractionResult(
            status="extract_error",
            backend="docling",
            backend_version=package_version("docling"),
            source_format=source_format,
            structure_origin="docling_tree",
            retryable=True,
            reason=_error_code("docling", exc),
        )
    if not any(block.text.strip() for block in blocks):
        return ExtractionResult(
            status="zero_text",
            backend="docling",
            backend_version=package_version("docling"),
            source_format=source_format,
            structure_origin="docling_tree",
            retryable=True,
            reason="docling_zero_text",
        )
    return ExtractionResult(
        status="indexed",
        backend="docling",
        backend_version=package_version("docling"),
        source_format=source_format,
        structure_origin="docling_tree",
        blocks=blocks,
        diagnostics={
            "options": DOCLING_OPTIONS_VERSION,
            "ocr": False,
            "source_range_basis": "canonical_extracted_text",
        },
    )


def _get_docling_converter(*, threads: int | None = None) -> Any:
    global _DOCLING_CONVERTER
    if _DOCLING_CONVERTER is not None:
        return _DOCLING_CONVERTER
    os.environ.setdefault("HF_HUB_OFFLINE", "1")
    os.environ.setdefault("TRANSFORMERS_OFFLINE", "1")
    from docling.datamodel.accelerator_options import (
        AcceleratorDevice,
        AcceleratorOptions,
    )
    from docling.datamodel.base_models import InputFormat
    from docling.datamodel.pipeline_options import PdfPipelineOptions
    from docling.document_converter import DocumentConverter, PdfFormatOption

    artifacts_path = docling_pdf_artifacts_path()
    allowed_formats = [InputFormat.DOCX, InputFormat.PPTX]
    format_options = {}
    if artifacts_path is not None:
        pipeline = PdfPipelineOptions()
        pipeline.do_ocr = False
        pipeline.do_table_structure = True
        pipeline.accelerator_options = AcceleratorOptions(
            num_threads=max(
                1,
                int(threads or os.getenv("OMP_NUM_THREADS", "1")),
            ),
            device=AcceleratorDevice.CPU,
        )
        pipeline.artifacts_path = artifacts_path
        allowed_formats.insert(0, InputFormat.PDF)
        format_options[InputFormat.PDF] = PdfFormatOption(
            pipeline_options=pipeline
        )
    _DOCLING_CONVERTER = DocumentConverter(
        allowed_formats=allowed_formats,
        format_options=format_options,
    )
    return _DOCLING_CONVERTER


def _docling_blocks(document: Any, source_format: str) -> Iterable[StructureBlock]:
    headings: list[str] = []
    raw: list[StructureBlock] = []
    ordinal = 0
    for item, tree_level in document.iterate_items():
        class_name = type(item).__name__.lower()
        label = str(getattr(item, "label", "")).lower()
        text = str(getattr(item, "text", "") or "").strip()
        if "sectionheader" in class_name or "section_header" in label:
            level = max(1, int(getattr(item, "level", tree_level or 1)))
            headings[level - 1 :] = [text or f"Section {level}"]
            continue
        if class_name == "titleitem" or label.endswith("title"):
            if text:
                headings[:] = [text]
            continue
        kind = "text"
        if "table" in class_name or label.endswith("table"):
            text = _docling_table_text(item, document)
            kind = "table"
        elif "listitem" in class_name or "list_item" in label:
            marker = str(getattr(item, "marker", "") or "-").strip()
            text = f"{marker} {text}".strip()
            kind = "list"
        elif "code" in class_name or label.endswith("code"):
            kind = "code"
        if not text:
            continue
        ordinal += 1
        page = _docling_page(item)
        title_parts = [part for part in headings if part]
        if not title_parts:
            title_parts = [f"Page {page}" if page else "Document"]
        title = " > ".join(title_parts)
        parent = _stable_id(
            "docling-parent",
            f"{title}:{page or 0}",
        )
        raw.append(
            StructureBlock(
                title=title,
                text=text,
                kind=kind,
                structure_id=f"docling-{ordinal:05d}",
                parent_section_id=parent,
                breadcrumb=" > ".join(headings),
                page=page if source_format != "pptx" else None,
                slide=page if source_format == "pptx" else None,
                preserve_layout=kind in {"table", "code", "list"},
            )
        )
    yield from _coalesce_docling_blocks(raw)


def _coalesce_docling_blocks(
    blocks: list[StructureBlock],
) -> Iterable[StructureBlock]:
    pending: list[StructureBlock] = []

    def flush() -> StructureBlock | None:
        nonlocal pending
        if not pending:
            return None
        first = pending[0]
        if len(pending) == 1:
            pending = []
            return first
        combined = StructureBlock(
            title=first.title,
            text="\n\n".join(block.text for block in pending),
            kind=(
                first.kind
                if all(block.kind == first.kind for block in pending)
                else "section"
            ),
            structure_id=f"{first.structure_id}--{pending[-1].structure_id}",
            parent_section_id=first.parent_section_id,
            breadcrumb=first.breadcrumb,
            page=first.page,
            slide=first.slide,
            preserve_layout=any(block.preserve_layout for block in pending),
        )
        pending = []
        return combined

    for block in blocks:
        same_parent = pending and (
            block.parent_section_id == pending[0].parent_section_id
            and block.kind != "table"
            and pending[0].kind != "table"
            and sum(len(value.text) for value in pending) + len(block.text)
            <= 6_000
        )
        if pending and not same_parent:
            value = flush()
            if value:
                yield value
        pending.append(block)
    value = flush()
    if value:
        yield value


def _with_canonical_offsets(
    blocks: Iterable[StructureBlock],
) -> Iterable[StructureBlock]:
    """Assign non-overlapping offsets in the canonical extracted text."""

    cursor = 0
    for block in blocks:
        yield replace(
            block,
            source_start=cursor,
            source_end=cursor + len(block.text),
        )
        cursor += len(block.text) + 2


def _docling_table_text(item: Any, document: Any) -> str:
    exporter = getattr(item, "export_to_dataframe", None)
    if callable(exporter):
        try:
            frame = exporter(doc=document)
        except TypeError:
            frame = exporter()
        try:
            return frame.to_csv(sep="\t", index=False).strip()
        except Exception:
            pass
    data = getattr(item, "data", None)
    grid = getattr(data, "grid", None)
    if grid:
        rows: list[str] = []
        for row in grid:
            values = [
                str(getattr(cell, "text", "") or "").strip()
                for cell in row
            ]
            rows.append("\t".join(values).rstrip())
        return "\n".join(row for row in rows if row.strip())
    return str(getattr(item, "text", "") or "").strip()


def _docling_page(item: Any) -> int | None:
    provenance = getattr(item, "prov", None) or []
    for value in provenance:
        page = getattr(value, "page_no", None)
        if isinstance(page, int):
            return page
    return None


def extract_legacy_office_structure(path: Path) -> ExtractionResult:
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            from pypdf import PdfReader

            blocks = [
                StructureBlock(
                    title=f"Page {index}",
                    text=page.extract_text() or "",
                    kind="page",
                    structure_id=f"legacy-page-{index:04d}",
                    parent_section_id=f"legacy-page-{index:04d}",
                    page=index,
                )
                for index, page in enumerate(
                    PdfReader(str(path)).pages, start=1
                )
            ]
            backend = "pypdf"
        elif ext == ".docx":
            from docx import Document

            doc = Document(str(path))
            parts = [
                paragraph.text for paragraph in doc.paragraphs
                if paragraph.text.strip()
            ]
            for table in doc.tables:
                for row in table.rows:
                    values = [cell.text.strip() for cell in row.cells]
                    if any(values):
                        parts.append("\t".join(values))
            blocks = [
                StructureBlock(
                    title=path.name,
                    text="\n".join(parts),
                    kind="document",
                    structure_id="legacy-docx-0001",
                    parent_section_id="legacy-docx-0001",
                )
            ]
            backend = "python-docx"
        elif ext == ".pptx":
            from pptx import Presentation

            blocks = []
            for index, slide in enumerate(
                Presentation(str(path)).slides, start=1
            ):
                parts = [
                    shape.text for shape in slide.shapes
                    if hasattr(shape, "text") and shape.text.strip()
                ]
                blocks.append(
                    StructureBlock(
                        title=f"Slide {index}",
                        text="\n".join(parts),
                        kind="slide",
                        structure_id=f"legacy-slide-{index:04d}",
                        parent_section_id=f"legacy-slide-{index:04d}",
                        slide=index,
                    )
                )
            backend = "python-pptx"
        else:
            return ExtractionResult(
                status="unsupported",
                backend="legacy",
                backend_version="",
                source_format=ext.lstrip("."),
                structure_origin="legacy",
                reason="legacy_format_unsupported",
            )
    except Exception as exc:
        return ExtractionResult(
            status="extract_error",
            backend="legacy",
            backend_version="",
            source_format=ext.lstrip("."),
            structure_origin="legacy",
            retryable=True,
            reason=_error_code("legacy_office", exc),
        )
    blocks = list(
        _with_canonical_offsets(
            block for block in blocks if block.text.strip()
        )
    )
    if not blocks:
        return ExtractionResult(
            status="zero_text",
            backend=backend,
            backend_version=_legacy_backend_version(backend),
            source_format=ext.lstrip("."),
            structure_origin="legacy",
            retryable=True,
            reason="legacy_zero_text",
        )
    return ExtractionResult(
        status="indexed",
        backend=backend,
        backend_version=_legacy_backend_version(backend),
        source_format=ext.lstrip("."),
        structure_origin="legacy",
        blocks=tuple(blocks),
        diagnostics={"source_range_basis": "canonical_extracted_text"},
    )


def _legacy_backend_version(backend: str) -> str:
    package = {
        "pypdf": "pypdf",
        "python-docx": "python-docx",
        "python-pptx": "python-pptx",
    }.get(backend, backend)
    return package_version(package)


def _normalize_line_endings(text: str) -> str:
    return text.replace("\r\n", "\n").replace("\r", "\n")


def _stable_id(prefix: str, value: str) -> str:
    digest = hashlib.sha256(value.encode("utf-8", errors="strict")).hexdigest()
    return f"{prefix}-{digest[:16]}"


def _error_code(stage: str, exc: BaseException) -> str:
    return f"{stage}_{type(exc).__name__}"[:160]
